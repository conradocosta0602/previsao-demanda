# -*- coding: utf-8 -*-
"""
Script para criar apresentação PowerPoint - Versão 4 Premium
Design elegante com diagramas elaborados e visual de consultoria
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Criar apresentação widescreen
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ============ PALETA DE CORES PREMIUM ============
AZUL_NAVY = RGBColor(15, 32, 65)
AZUL_ROYAL = RGBColor(30, 58, 95)
AZUL_ACCENT = RGBColor(0, 120, 200)
AZUL_LIGHT = RGBColor(70, 130, 180)

DOURADO = RGBColor(218, 165, 32)
AMBAR_LIGHT = RGBColor(255, 193, 7)

BRANCO = RGBColor(255, 255, 255)
CINZA_CLARO = RGBColor(240, 242, 245)
CINZA_MEDIO = RGBColor(160, 170, 180)
CINZA_ESCURO = RGBColor(80, 90, 100)

VERDE_SUCESSO = RGBColor(40, 167, 69)
VERMELHO_ALERTA = RGBColor(220, 53, 69)
LARANJA_ATENCAO = RGBColor(255, 140, 0)
ROXO = RGBColor(128, 90, 180)


def add_premium_background(slide, style="navy"):
    """Adiciona fundo premium"""
    if style == "navy":
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = AZUL_NAVY
        bg.line.fill.background()

        accent = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2), Inches(6), Inches(6))
        accent.fill.solid()
        accent.fill.fore_color.rgb = AZUL_ROYAL
        accent.line.fill.background()

    elif style == "light":
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = CINZA_CLARO
        bg.line.fill.background()

    elif style == "gradient":
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = BRANCO
        bg.line.fill.background()

        side = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.4), prs.slide_height)
        side.fill.solid()
        side.fill.fore_color.rgb = AZUL_NAVY
        side.line.fill.background()


def add_elegant_header(slide, title, subtitle="", with_line=True):
    """Adiciona header elegante"""
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title.upper()
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = AZUL_NAVY
    p.font.name = "Calibri Light"

    if subtitle:
        p = tf.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(14)
        p.font.color.rgb = CINZA_MEDIO
        p.space_before = Pt(5)

    if with_line:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.35), Inches(1.5), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = DOURADO
        line.line.fill.background()


def add_stat_card(slide, x, y, number, label, sublabel="", accent_color=None, width=2.5, height=1.8):
    """Card de estatística elegante"""
    if accent_color is None:
        accent_color = AZUL_ACCENT

    shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x+0.05), Inches(y+0.05), Inches(width), Inches(height))
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = RGBColor(220, 220, 220)
    shadow.line.fill.background()

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = BRANCO
    card.line.fill.background()

    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(0.08))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent_color
    accent_bar.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(x), Inches(y + 0.25), Inches(width), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = accent_color
    p.alignment = PP_ALIGN.CENTER

    txBox = slide.shapes.add_textbox(Inches(x), Inches(y + 1.1), Inches(width), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = label.upper()
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = CINZA_ESCURO
    p.alignment = PP_ALIGN.CENTER

    if sublabel:
        txBox = slide.shapes.add_textbox(Inches(x), Inches(y + 1.4), Inches(width), Inches(0.3))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = sublabel
        p.font.size = Pt(9)
        p.font.color.rgb = CINZA_MEDIO
        p.alignment = PP_ALIGN.CENTER


def add_elegant_box(slide, x, y, width, height, title, items, icon_letter, color):
    """Box elegante com ícone circular"""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = BRANCO
    box.line.color.rgb = RGBColor(230, 230, 230)
    box.line.width = Pt(1)

    icon_size = 0.6
    icon = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.2), Inches(y + 0.2), Inches(icon_size), Inches(icon_size))
    icon.fill.solid()
    icon.fill.fore_color.rgb = color
    icon.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.3), Inches(icon_size), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = icon_letter
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER

    txBox = slide.shapes.add_textbox(Inches(x + 1), Inches(y + 0.25), Inches(width - 1.2), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = AZUL_NAVY

    txBox = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.9), Inches(width - 0.6), Inches(height - 1.1))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(12)
        p.font.color.rgb = CINZA_ESCURO
        p.space_after = Pt(4)


def add_hexagon(slide, x, y, size, color, text, subtext=""):
    """Adiciona hexágono estilizado"""
    hex_shape = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(x), Inches(y), Inches(size), Inches(size * 0.9))
    hex_shape.fill.solid()
    hex_shape.fill.fore_color.rgb = color
    hex_shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(x), Inches(y + size * 0.25), Inches(size), Inches(size * 0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER

    if subtext:
        txBox = slide.shapes.add_textbox(Inches(x - 0.2), Inches(y + size), Inches(size + 0.4), Inches(0.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = subtext
        p.font.size = Pt(9)
        p.font.color.rgb = CINZA_ESCURO
        p.alignment = PP_ALIGN.CENTER


def add_flow_arrow(slide, x, y, width, color=None):
    """Seta de fluxo elegante"""
    if color is None:
        color = DOURADO
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(width), Inches(0.25))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()


def add_chevron_arrow(slide, x, y, width, height, color):
    """Chevron arrow para fluxos"""
    chevron = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(width), Inches(height))
    chevron.fill.solid()
    chevron.fill.fore_color.rgb = color
    chevron.line.fill.background()


# ============ CRIAR SLIDES ============

# ======== SLIDE 1 - CAPA ========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_premium_background(slide, "navy")

geo1 = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, Inches(9), Inches(-1), Inches(6), Inches(10))
geo1.fill.solid()
geo1.fill.fore_color.rgb = AZUL_ROYAL
geo1.line.fill.background()

gold_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.8), Inches(1.2), Inches(0.06))
gold_line.fill.solid()
gold_line.fill.fore_color.rgb = DOURADO
gold_line.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(3), Inches(8), Inches(1.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "PREVISAO DE"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = BRANCO
p.font.name = "Calibri Light"

p = tf.add_paragraph()
p.text = "DEMANDA"
p.font.size = Pt(60)
p.font.bold = True
p.font.color.rgb = DOURADO
p.font.name = "Calibri Light"
p.space_before = Pt(0)

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(7), Inches(0.8))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Sistema Inteligente de Gestao de Pedidos"
p.font.size = Pt(18)
p.font.color.rgb = CINZA_CLARO

p = tf.add_paragraph()
p.text = "Proposta de Implantacao Piloto"
p.font.size = Pt(14)
p.font.color.rgb = CINZA_MEDIO
p.space_before = Pt(10)

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(3), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Janeiro 2026"
p.font.size = Pt(12)
p.font.color.rgb = CINZA_MEDIO


# ======== SLIDE 2 - O PROBLEMA (Diagrama Elaborado) ========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_premium_background(slide, "gradient")
add_elegant_header(slide, "O Cenario Atual", "Impactos na operacao e no resultado")

# Diagrama de funil invertido mostrando problemas cascateando
# Problema raiz no topo
root_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(1.7), Inches(4.3), Inches(0.9))
root_box.fill.solid()
root_box.fill.fore_color.rgb = VERMELHO_ALERTA
root_box.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(4.5), Inches(1.85), Inches(4.3), Inches(0.6))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "GESTAO MANUAL DE PEDIDOS"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

# Setas descendo
arrow_down1 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(5.5), Inches(2.65), Inches(0.4), Inches(0.5))
arrow_down1.fill.solid()
arrow_down1.fill.fore_color.rgb = VERMELHO_ALERTA
arrow_down1.line.fill.background()

arrow_down2 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(7.4), Inches(2.65), Inches(0.4), Inches(0.5))
arrow_down2.fill.solid()
arrow_down2.fill.fore_color.rgb = VERMELHO_ALERTA
arrow_down2.line.fill.background()

# Segunda linha - dois problemas
prob1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.8), Inches(3.2), Inches(3.2), Inches(0.8))
prob1.fill.solid()
prob1.fill.fore_color.rgb = RGBColor(255, 200, 200)
prob1.line.color.rgb = VERMELHO_ALERTA
prob1.line.width = Pt(2)

txBox = slide.shapes.add_textbox(Inches(2.8), Inches(3.35), Inches(3.2), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Sem previsao estruturada"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = AZUL_NAVY
p.alignment = PP_ALIGN.CENTER

prob2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.3), Inches(3.2), Inches(3.2), Inches(0.8))
prob2.fill.solid()
prob2.fill.fore_color.rgb = RGBColor(255, 200, 200)
prob2.line.color.rgb = VERMELHO_ALERTA
prob2.line.width = Pt(2)

txBox = slide.shapes.add_textbox(Inches(7.3), Inches(3.35), Inches(3.2), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Mesma regra para tudo"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = AZUL_NAVY
p.alignment = PP_ALIGN.CENTER

# Setas para terceira linha
for x in [3.9, 5.9, 8.4, 10.4]:
    if x < 7:
        arr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x - 0.5), Inches(4.05), Inches(0.3), Inches(0.4))
    else:
        arr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x - 0.5), Inches(4.05), Inches(0.3), Inches(0.4))
    arr.fill.solid()
    arr.fill.fore_color.rgb = RGBColor(255, 150, 150)
    arr.line.fill.background()

# Terceira linha - impactos operacionais
impactos = [
    (0.8, "RUPTURAS", "Falta de produto\nna gondola", VERMELHO_ALERTA),
    (3.3, "EXCESSO", "Estoque parado\nCapital imobilizado", LARANJA_ATENCAO),
    (5.8, "RETRABALHO", "Ajustes manuais\nconstantes", LARANJA_ATENCAO),
    (8.3, "INEFICIENCIA", "Horas em calculos\nrepetitivos", LARANJA_ATENCAO),
    (10.8, "SEM VISAO", "Decisao por\nintuicao", VERMELHO_ALERTA)
]

for x, titulo, desc, color in impactos:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(4.5), Inches(2.3), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = BRANCO
    box.line.color.rgb = color
    box.line.width = Pt(2)

    # Barra colorida no topo
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(4.5), Inches(2.3), Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(x), Inches(4.55), Inches(2.3), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER

    txBox = slide.shapes.add_textbox(Inches(x), Inches(5), Inches(2.3), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(10)
    p.font.color.rgb = CINZA_ESCURO
    p.alignment = PP_ALIGN.CENTER

# Box de impacto financeiro
impact_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(6.3), Inches(9.3), Inches(0.7))
impact_box.fill.solid()
impact_box.fill.fore_color.rgb = AZUL_NAVY
impact_box.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(2), Inches(6.45), Inches(9.3), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "RESULTADO:  Perda de vendas  |  Custos elevados  |  Baixa produtividade  |  Cliente insatisfeito"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER


# ======== SLIDE 3 - A SOLUÇÃO (Fluxo Elaborado com Pipeline) ========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_premium_background(slide, "gradient")
add_elegant_header(slide, "A Transformacao", "Pipeline inteligente de dados para decisao")

# Pipeline visual com chevrons
# Etapa 1 - COLETA
add_chevron_arrow(slide, 0.5, 2.5, 2.8, 1.2, AZUL_ACCENT)
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.75), Inches(2.4), Inches(0.7))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "COLETA"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

# Etapa 2 - ANALISE
add_chevron_arrow(slide, 2.9, 2.5, 2.8, 1.2, RGBColor(0, 100, 180))
txBox = slide.shapes.add_textbox(Inches(2.9), Inches(2.75), Inches(2.4), Inches(0.7))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "ANALISE"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

# Etapa 3 - PREVISAO
add_chevron_arrow(slide, 5.3, 2.5, 2.8, 1.2, RGBColor(0, 80, 160))
txBox = slide.shapes.add_textbox(Inches(5.3), Inches(2.75), Inches(2.4), Inches(0.7))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "PREVISAO"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

# Etapa 4 - ACAO
add_chevron_arrow(slide, 7.7, 2.5, 2.8, 1.2, AZUL_NAVY)
txBox = slide.shapes.add_textbox(Inches(7.7), Inches(2.75), Inches(2.4), Inches(0.7))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "ACAO"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

# Etapa 5 - RESULTADO
add_chevron_arrow(slide, 10.1, 2.5, 2.8, 1.2, VERDE_SUCESSO)
txBox = slide.shapes.add_textbox(Inches(10.1), Inches(2.75), Inches(2.4), Inches(0.7))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "RESULTADO"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

# Detalhes abaixo de cada etapa
detalhes = [
    (0.5, ["10M+ registros", "3 anos historico", "Todas as lojas"]),
    (2.9, ["Sazonalidade", "Tendencias", "Padroes de venda"]),
    (5.3, ["Forecast 12 meses", "Por SKU e loja", "Machine Learning"]),
    (7.7, ["Pedidos otimizados", "Alertas de ruptura", "Balanceamento"]),
    (10.1, ["Menos ruptura", "Estoque ideal", "Mais vendas"])
]

for x, items in detalhes:
    y = 4
    for item in items:
        txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(2.6), Inches(0.35))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"  {item}"
        p.font.size = Pt(11)
        p.font.color.rgb = CINZA_ESCURO
        y += 0.4

# Diagrama circular central - O Sistema
center_x, center_y = 5.4, 5.5
system_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(center_x), Inches(center_y), Inches(2.5), Inches(1.3))
system_circle.fill.solid()
system_circle.fill.fore_color.rgb = AZUL_NAVY
system_circle.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(center_x), Inches(center_y + 0.25), Inches(2.5), Inches(0.8))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "SISTEMA DE"
p.font.size = Pt(12)
p.font.color.rgb = CINZA_CLARO
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "DEMANDA"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = DOURADO
p.alignment = PP_ALIGN.CENTER

# Setas conectando ao sistema
for i, x_pos in enumerate([1.5, 3.8, 7, 9.3]):
    conn = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x_pos + 0.8), Inches(5.1), Inches(0.3), Inches(0.35))
    conn.fill.solid()
    conn.fill.fore_color.rgb = DOURADO
    conn.line.fill.background()


# ======== SLIDE 4 - MÓDULOS DO SISTEMA (Atualizado) ========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_premium_background(slide, "gradient")
add_elegant_header(slide, "Funcionalidades", "Quatro pilares integrados")

modulos = [
    (0.8, 1.8, "Previsao de Demanda", "P", AZUL_ACCENT,
     ["Forecast de 12 meses", "Analise por produto e loja", "Deteccao de sazonalidade", "Ajuste de tendencias"]),
    (6.9, 1.8, "Gestao de Pedidos", "G", VERDE_SUCESSO,
     ["Sugestao automatica", "Arredondamento por embalagem", "Calendario de entregas", "Gestao de fornecedores"]),
    (0.8, 4.4, "Controle de Estoque", "E", LARANJA_ATENCAO,
     ["Posicao em tempo real", "Alertas de ruptura", "Balanceamento entre lojas", "Indicadores de giro"]),
    (6.9, 4.4, "Cenarios e Promocoes", "C", ROXO,
     ["Simulacao de cenarios", "Eventos promocionais", "Impacto de acoes", "Ajustes de demanda"])
]

for x, y, title, icon, color, items in modulos:
    add_elegant_box(slide, x, y, 6, 2.4, title, items, icon, color)


# ======== SLIDE 5 - ANTES vs DEPOIS (Visual Elaborado) ========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_premium_background(slide, "gradient")
add_elegant_header(slide, "A Evolucao do Processo", "De reativo para proativo")

# Timeline horizontal mostrando evolução
# Linha central
timeline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.4), Inches(11.3), Inches(0.08))
timeline.fill.solid()
timeline.fill.fore_color.rgb = DOURADO
timeline.line.fill.background()

# LADO ESQUERDO - ANTES
before_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), Inches(2.9), Inches(1), Inches(1))
before_circle.fill.solid()
before_circle.fill.fore_color.rgb = VERMELHO_ALERTA
before_circle.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(1.5), Inches(3.2), Inches(1), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "HOJE"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

# Cards do ANTES
antes_items = [
    ("Manual", "Calculos em planilha"),
    ("Padronizado", "Mesma regra para todos"),
    ("Reativo", "Corre atras do problema"),
]

x_antes = 0.8
for titulo, desc in antes_items:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_antes), Inches(4.1), Inches(1.9), Inches(1.4))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(255, 240, 240)
    box.line.color.rgb = VERMELHO_ALERTA
    box.line.width = Pt(1)

    txBox = slide.shapes.add_textbox(Inches(x_antes), Inches(4.25), Inches(1.9), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = VERMELHO_ALERTA
    p.alignment = PP_ALIGN.CENTER

    txBox = slide.shapes.add_textbox(Inches(x_antes), Inches(4.7), Inches(1.9), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(10)
    p.font.color.rgb = CINZA_ESCURO
    p.alignment = PP_ALIGN.CENTER

    x_antes += 2.1

# CENTRO - Transformação
transform_box = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.9), Inches(2.6), Inches(1.5), Inches(1.5))
transform_box.fill.solid()
transform_box.fill.fore_color.rgb = DOURADO
transform_box.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(5.9), Inches(3.05), Inches(1.5), Inches(0.7))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "TRANS"
p.font.size = Pt(10)
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "FORMACAO"
p.font.size = Pt(10)
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

# LADO DIREITO - DEPOIS
after_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.8), Inches(2.9), Inches(1), Inches(1))
after_circle.fill.solid()
after_circle.fill.fore_color.rgb = VERDE_SUCESSO
after_circle.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(10.8), Inches(3.2), Inches(1), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "NOVO"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

# Cards do DEPOIS
depois_items = [
    ("Automatico", "Sistema calcula tudo"),
    ("Personalizado", "Por item e loja"),
    ("Proativo", "Antecipa problemas"),
]

x_depois = 7.5
for titulo, desc in depois_items:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_depois), Inches(4.1), Inches(1.9), Inches(1.4))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(235, 255, 235)
    box.line.color.rgb = VERDE_SUCESSO
    box.line.width = Pt(1)

    txBox = slide.shapes.add_textbox(Inches(x_depois), Inches(4.25), Inches(1.9), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = VERDE_SUCESSO
    p.alignment = PP_ALIGN.CENTER

    txBox = slide.shapes.add_textbox(Inches(x_depois), Inches(4.7), Inches(1.9), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(10)
    p.font.color.rgb = CINZA_ESCURO
    p.alignment = PP_ALIGN.CENTER

    x_depois += 2.1

# Box de destaque inferior
highlight_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(5.8), Inches(10.3), Inches(1.2))
highlight_box.fill.solid()
highlight_box.fill.fore_color.rgb = AZUL_NAVY
highlight_box.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(1.5), Inches(6), Inches(10.3), Inches(0.8))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "RESULTADO DA TRANSFORMACAO"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = DOURADO
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "Decisoes baseadas em dados  |  Pedidos precisos  |  Equipe focada em estrategia"
p.font.size = Pt(12)
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER


# ======== SLIDE 6 - PROPOSTA DO PILOTO ========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_premium_background(slide, "gradient")
add_elegant_header(slide, "Proposta de Piloto", "Escopo controlado, resultados rapidos")

add_stat_card(slide, 0.8, 1.8, "4", "FORNECEDORES", "Selecionados", AZUL_ACCENT, 2.8, 1.6)
add_stat_card(slide, 4, 1.8, "701", "PRODUTOS", "No escopo", VERDE_SUCESSO, 2.8, 1.6)
add_stat_card(slide, 7.2, 1.8, "10M+", "REGISTROS", "Historico", LARANJA_ATENCAO, 2.8, 1.6)
add_stat_card(slide, 10.4, 1.8, "3 anos", "DADOS", "Disponiveis", ROXO, 2.8, 1.6)

fornecedores = [
    ("FAME", "45 SKUs", "Ferragens", AZUL_ACCENT),
    ("ZAGONEL", "75 SKUs", "Chuveiros/Duchas", VERDE_SUCESSO),
    ("SILVANA", "311 SKUs", "Fechaduras", LARANJA_ATENCAO),
    ("STAM", "270 SKUs", "Fechaduras", ROXO)
]

x_pos = 0.8
for nome, skus, cat, color in fornecedores:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(4), Inches(3), Inches(2.2))
    card.fill.solid()
    card.fill.fore_color.rgb = BRANCO
    card.line.color.rgb = color
    card.line.width = Pt(2)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x_pos), Inches(4), Inches(0.12), Inches(2.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(x_pos + 0.3), Inches(4.2), Inches(2.5), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = nome
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = AZUL_NAVY

    txBox = slide.shapes.add_textbox(Inches(x_pos + 0.3), Inches(4.8), Inches(2.5), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = skus
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = color

    txBox = slide.shapes.add_textbox(Inches(x_pos + 0.3), Inches(5.4), Inches(2.5), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = cat
    p.font.size = Pt(12)
    p.font.color.rgb = CINZA_MEDIO

    x_pos += 3.2

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(12), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Todas as lojas da rede   |   Sistema pronto para uso   |   Dados ja carregados e validados"
p.font.size = Pt(12)
p.font.color.rgb = CINZA_MEDIO
p.alignment = PP_ALIGN.CENTER


# ======== SLIDE 7 - BENEFÍCIOS (Números Suavizados) ========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_premium_background(slide, "gradient")
add_elegant_header(slide, "Beneficios Esperados", "Metas conservadoras com alto potencial")

# Números suavizados (exceto produtividade)
beneficios = [
    ("-15%", "REDUCAO DE RUPTURA", "Mais produtos disponiveis para venda", VERDE_SUCESSO),
    ("+20%", "ACURACIDADE", "Pedidos mais assertivos", AZUL_ACCENT),
    ("+15%", "EFICIENCIA DE ESTOQUE", "Melhor balanceamento", LARANJA_ATENCAO),
    ("+50%", "PRODUTIVIDADE", "Menos tempo em calculos manuais", ROXO)
]

x_pos = 0.8
for valor, titulo, desc, color in beneficios:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(1.8), Inches(3), Inches(3.8))
    card.fill.solid()
    card.fill.fore_color.rgb = BRANCO
    card.line.color.rgb = RGBColor(230, 230, 230)
    card.line.width = Pt(1)

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x_pos + 0.7), Inches(2.1), Inches(1.6), Inches(1.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(x_pos + 0.7), Inches(2.5), Inches(1.6), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = valor
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER

    txBox = slide.shapes.add_textbox(Inches(x_pos), Inches(3.9), Inches(3), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = AZUL_NAVY
    p.alignment = PP_ALIGN.CENTER

    txBox = slide.shapes.add_textbox(Inches(x_pos + 0.2), Inches(4.4), Inches(2.6), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(10)
    p.font.color.rgb = CINZA_MEDIO
    p.alignment = PP_ALIGN.CENTER

    x_pos += 3.15

impact_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6), Inches(12.5), Inches(0.9))
impact_box.fill.solid()
impact_box.fill.fore_color.rgb = AZUL_NAVY
impact_box.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(12.5), Inches(0.6))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "METAS CONSERVADORAS  |  Potencial de superacao significativo com ajustes e aprendizado continuo"
p.font.size = Pt(14)
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER


# ======== SLIDE 8 - BENEFÍCIOS DETALHADOS ========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_premium_background(slide, "gradient")
add_elegant_header(slide, "Como Geramos Valor", "Cadeia de beneficios integrados")

# Diagrama de fluxo mostrando como cada benefício gera valor
# Fluxo: Previsao -> Pedidos -> Estoque -> Vendas -> Resultado

# Box 1 - Previsão
box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2), Inches(2.3), Inches(2.2))
box1.fill.solid()
box1.fill.fore_color.rgb = AZUL_ACCENT
box1.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.15), Inches(2.3), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "PREVISAO"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

txBox = slide.shapes.add_textbox(Inches(0.6), Inches(2.6), Inches(2.1), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Antecipar demanda\ncom precisao"
p.font.size = Pt(11)
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

# Seta
add_flow_arrow(slide, 2.9, 2.9, 0.5)

# Box 2 - Pedidos
box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(2), Inches(2.3), Inches(2.2))
box2.fill.solid()
box2.fill.fore_color.rgb = VERDE_SUCESSO
box2.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(3.5), Inches(2.15), Inches(2.3), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "PEDIDOS"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

txBox = slide.shapes.add_textbox(Inches(3.6), Inches(2.6), Inches(2.1), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Quantidade certa\nno momento certo"
p.font.size = Pt(11)
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

# Seta
add_flow_arrow(slide, 5.9, 2.9, 0.5)

# Box 3 - Estoque
box3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(2), Inches(2.3), Inches(2.2))
box3.fill.solid()
box3.fill.fore_color.rgb = LARANJA_ATENCAO
box3.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(6.5), Inches(2.15), Inches(2.3), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "ESTOQUE"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

txBox = slide.shapes.add_textbox(Inches(6.6), Inches(2.6), Inches(2.1), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Nivel ideal em\ncada loja"
p.font.size = Pt(11)
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

# Seta
add_flow_arrow(slide, 8.9, 2.9, 0.5)

# Box 4 - Resultado
box4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.5), Inches(2), Inches(3.3), Inches(2.2))
box4.fill.solid()
box4.fill.fore_color.rgb = AZUL_NAVY
box4.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(9.5), Inches(2.15), Inches(3.3), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "RESULTADO"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = DOURADO
p.alignment = PP_ALIGN.CENTER

txBox = slide.shapes.add_textbox(Inches(9.6), Inches(2.6), Inches(3.1), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "+ Vendas\n- Custos\n+ Produtividade"
p.font.size = Pt(12)
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

# Cards de impacto detalhado abaixo
impactos_det = [
    ("OPERACIONAL", ["Menos horas em calculos", "Processo padronizado", "Decisao mais rapida"], AZUL_ACCENT),
    ("FINANCEIRO", ["Capital de giro otimizado", "Menor custo logistico", "Fluxo de caixa melhor"], VERDE_SUCESSO),
    ("COMERCIAL", ["Menos ruptura na loja", "Cliente encontra produto", "Mais vendas realizadas"], LARANJA_ATENCAO),
    ("ESTRATEGICO", ["Dados para negociacao", "Visibilidade de demanda", "Planejamento de acoes"], ROXO)
]

x_imp = 0.5
for titulo, items, color in impactos_det:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_imp), Inches(4.5), Inches(3.1), Inches(2.3))
    box.fill.solid()
    box.fill.fore_color.rgb = BRANCO
    box.line.color.rgb = color
    box.line.width = Pt(2)

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x_imp), Inches(4.5), Inches(3.1), Inches(0.5))
    header.fill.solid()
    header.fill.fore_color.rgb = color
    header.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(x_imp), Inches(4.58), Inches(3.1), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER

    # Items
    txBox = slide.shapes.add_textbox(Inches(x_imp + 0.2), Inches(5.1), Inches(2.7), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(10)
        p.font.color.rgb = CINZA_ESCURO

    x_imp += 3.2


# ======== SLIDE 9 - ROADMAP (Elaborado com Fases) ========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_premium_background(slide, "gradient")
add_elegant_header(slide, "Roadmap de Expansao", "Crescimento gradual e controlado")

# Visual de montanha/crescimento
# Base
base = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6), Inches(12.3), Inches(0.15))
base.fill.solid()
base.fill.fore_color.rgb = CINZA_MEDIO
base.line.fill.background()

# Fase 1 - Piloto (destacado)
fase1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.5), Inches(2.8), Inches(2.5))
fase1.fill.solid()
fase1.fill.fore_color.rgb = AZUL_ACCENT
fase1.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(3.7), Inches(2.8), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "FASE 1"
p.font.size = Pt(11)
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(4.1), Inches(2.8), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "PILOTO"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

txBox = slide.shapes.add_textbox(Inches(0.9), Inches(4.7), Inches(2.6), Inches(1.1))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "4 fornecedores\n701 SKUs\nValidacao do modelo"
p.font.size = Pt(11)
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

# Seta
arr1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.7), Inches(4.5), Inches(0.5), Inches(0.3))
arr1.fill.solid()
arr1.fill.fore_color.rgb = DOURADO
arr1.line.fill.background()

# Fase 2
fase2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.3), Inches(2.8), Inches(2.8), Inches(3.2))
fase2.fill.solid()
fase2.fill.fore_color.rgb = RGBColor(0, 100, 180)
fase2.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(4.3), Inches(3), Inches(2.8), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "FASE 2"
p.font.size = Pt(11)
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

txBox = slide.shapes.add_textbox(Inches(4.3), Inches(3.4), Inches(2.8), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "EXPANSAO I"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

txBox = slide.shapes.add_textbox(Inches(4.4), Inches(4), Inches(2.6), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "+8 fornecedores\n~2.000 SKUs\nConsolidacao"
p.font.size = Pt(11)
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

# Seta
arr2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(7.2), Inches(4), Inches(0.5), Inches(0.3))
arr2.fill.solid()
arr2.fill.fore_color.rgb = DOURADO
arr2.line.fill.background()

# Fase 3
fase3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(2.1), Inches(2.8), Inches(3.9))
fase3.fill.solid()
fase3.fill.fore_color.rgb = RGBColor(0, 80, 160)
fase3.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(7.8), Inches(2.3), Inches(2.8), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "FASE 3"
p.font.size = Pt(11)
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

txBox = slide.shapes.add_textbox(Inches(7.8), Inches(2.7), Inches(2.8), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "EXPANSAO II"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

txBox = slide.shapes.add_textbox(Inches(7.9), Inches(3.3), Inches(2.6), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "+11 fornecedores\n~5.000 SKUs\nEscala"
p.font.size = Pt(11)
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

# Seta
arr3 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(10.7), Inches(3.5), Inches(0.5), Inches(0.3))
arr3.fill.solid()
arr3.fill.fore_color.rgb = DOURADO
arr3.line.fill.background()

# Fase 4 - Completo
fase4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.3), Inches(1.5), Inches(1.7), Inches(4.5))
fase4.fill.solid()
fase4.fill.fore_color.rgb = VERDE_SUCESSO
fase4.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(11.3), Inches(1.7), Inches(1.7), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "FASE 4"
p.font.size = Pt(11)
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

txBox = slide.shapes.add_textbox(Inches(11.3), Inches(2.1), Inches(1.7), Inches(0.8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "TODAS CATEGORIAS"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

txBox = slide.shapes.add_textbox(Inches(11.3), Inches(3), Inches(1.7), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "23 cat.\nTodos\nSKUs"
p.font.size = Pt(11)
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

# Critérios
criteria_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.3), Inches(12.5), Inches(0.9))
criteria_box.fill.solid()
criteria_box.fill.fore_color.rgb = BRANCO
criteria_box.line.color.rgb = DOURADO
criteria_box.line.width = Pt(2)

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(12.5), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "CRITERIOS DE AVANCO:  Resultados positivos  |  Dados disponiveis  |  Capacidade operacional"
p.font.size = Pt(12)
p.font.color.rgb = AZUL_NAVY
p.alignment = PP_ALIGN.CENTER


# ======== SLIDE 10 - CRONOGRAMA EM MESES ========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_premium_background(slide, "gradient")
add_elegant_header(slide, "Cronograma do Piloto", "Fases em meses")

# Gantt chart em meses
fases_gantt = [
    ("Setup e Configuracao", 0, 3, AZUL_ACCENT, ["Configuracao do ambiente", "Integracao de dados", "Validacao inicial"]),
    ("Treinamento da Equipe", 3, 1, VERDE_SUCESSO, ["Capacitacao usuarios", "Documentacao", "Hands-on"]),
    ("Operacao Assistida", 4, 2, LARANJA_ATENCAO, ["Acompanhamento proximo", "Ajustes finos", "Suporte dedicado"]),
    ("Operacao Autonoma", 6, 2, ROXO, ["Equipe autonoma", "Monitoramento", "Avaliacao final"])
]

# Header com meses
meses = ["M1", "M2", "M3", "M4", "M5", "M6", "M7", "M8"]
for i, mes in enumerate(meses):
    x = 5 + i * 1
    txBox = slide.shapes.add_textbox(Inches(x), Inches(1.7), Inches(1), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = mes
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = AZUL_NAVY
    p.alignment = PP_ALIGN.CENTER

    # Linha vertical
    vline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.5), Inches(2), Inches(0.01), Inches(4))
    vline.fill.solid()
    vline.fill.fore_color.rgb = RGBColor(230, 230, 230)
    vline.line.fill.background()

y_pos = 2.2
for fase, start, duration, color, details in fases_gantt:
    # Label da fase
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos + 0.15), Inches(4), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = fase
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = AZUL_NAVY

    # Barra do Gantt
    bar_x = 5 + start * 1
    bar_width = duration * 1

    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(bar_x), Inches(y_pos), Inches(bar_width), Inches(0.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    # Duração na barra
    txBox = slide.shapes.add_textbox(Inches(bar_x), Inches(y_pos + 0.15), Inches(bar_width), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    if duration == 1:
        p.text = "1 mes"
    else:
        p.text = f"{duration} meses"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER

    # Detalhes abaixo
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos + 0.6), Inches(4), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = " | ".join(details)
    p.font.size = Pt(9)
    p.font.color.rgb = CINZA_MEDIO

    y_pos += 1.1

# Marcos
marcos_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.3), Inches(12.5), Inches(0.9))
marcos_box.fill.solid()
marcos_box.fill.fore_color.rgb = AZUL_NAVY
marcos_box.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(12.5), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "MARCOS:  Mes 3 - Sistema validado  |  Mes 5 - Primeiro ciclo completo  |  Mes 8 - Avaliacao e decisao"
p.font.size = Pt(12)
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER


# ======== SLIDE 11 - POR QUE AGORA (Elaborado) ========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_premium_background(slide, "gradient")
add_elegant_header(slide, "Por Que Agora?", "Convergencia de fatores favoraveis")

# Diagrama de convergência - 5 elementos convergindo para centro
center_x, center_y = 5.9, 4.2
center_size = 1.8

# Centro - Oportunidade
center = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(center_x), Inches(center_y), Inches(center_size), Inches(center_size))
center.fill.solid()
center.fill.fore_color.rgb = DOURADO
center.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(center_x), Inches(center_y + 0.5), Inches(center_size), Inches(0.8))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "MOMENTO"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "IDEAL"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

# 5 fatores ao redor
fatores = [
    (2.5, 1.8, "DADOS PRONTOS", "10M+ registros\ncarregados e validados", AZUL_ACCENT),
    (9, 1.8, "SISTEMA FUNCIONAL", "Pronto para uso\nimediato", VERDE_SUCESSO),
    (0.8, 4.5, "BAIXO INVESTIMENTO", "Infraestrutura\nja disponivel", LARANJA_ATENCAO),
    (10.5, 4.5, "ALTO RETORNO", "Potencial de ganho\ncomprovado", ROXO),
    (5.9, 1.5, "RISCO CONTROLADO", "Piloto permite\najustes rapidos", RGBColor(100, 149, 237))
]

for x, y, titulo, desc, color in fatores:
    # Box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.5), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = BRANCO
    box.line.color.rgb = color
    box.line.width = Pt(2)

    # Barra colorida
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(2.5), Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(x), Inches(y + 0.05), Inches(2.5), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER

    txBox = slide.shapes.add_textbox(Inches(x), Inches(y + 0.5), Inches(2.5), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(10)
    p.font.color.rgb = CINZA_ESCURO
    p.alignment = PP_ALIGN.CENTER

# CTA
cta_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(6.3), Inches(7.3), Inches(0.9))
cta_box.fill.solid()
cta_box.fill.fore_color.rgb = AZUL_NAVY
cta_box.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(3), Inches(6.45), Inches(7.3), Inches(0.6))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "VAMOS COMECAR?  |  O piloto pode iniciar apos aprovacao"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER


# ======== SLIDE 12 - ENCERRAMENTO ========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_premium_background(slide, "navy")

geo = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, Inches(-2), Inches(0), Inches(8), Inches(10))
geo.fill.solid()
geo.fill.fore_color.rgb = AZUL_ROYAL
geo.line.fill.background()

gold_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.8), Inches(3.2), Inches(1.8), Inches(0.06))
gold_line.fill.solid()
gold_line.fill.fore_color.rgb = DOURADO
gold_line.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(0), Inches(3.4), Inches(13.333), Inches(1.2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "OBRIGADO"
p.font.size = Pt(64)
p.font.bold = True
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

txBox = slide.shapes.add_textbox(Inches(0), Inches(4.8), Inches(13.333), Inches(0.6))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Duvidas e proximos passos"
p.font.size = Pt(20)
p.font.color.rgb = CINZA_CLARO
p.alignment = PP_ALIGN.CENTER


# ============ SALVAR ============
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'Apresentacao_Previsao_Demanda_v4.pptx')
prs.save(output_path)
print(f'Apresentacao v4 salva em: {output_path}')
print('Versao 4 - Design Premium com Diagramas Elaborados')
