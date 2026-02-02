# -*- coding: utf-8 -*-
"""
Script para criar apresentação PowerPoint - Versão Premium/Consultoria
Design elegante, comercial e impactante para apresentação executiva
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import nsmap
import os

# Criar apresentação widescreen
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ============ PALETA DE CORES PREMIUM ============
# Azuis sofisticados
AZUL_NAVY = RGBColor(15, 32, 65)        # Fundo principal
AZUL_ROYAL = RGBColor(30, 58, 95)       # Fundo secundário
AZUL_ACCENT = RGBColor(0, 120, 200)     # Destaques
AZUL_LIGHT = RGBColor(70, 130, 180)     # Texto secundário

# Dourado/Âmbar (premium feel)
DOURADO = RGBColor(218, 165, 32)        # Accent premium
AMBAR_LIGHT = RGBColor(255, 193, 7)     # Highlights

# Neutros
BRANCO = RGBColor(255, 255, 255)
CINZA_CLARO = RGBColor(240, 242, 245)
CINZA_MEDIO = RGBColor(160, 170, 180)
CINZA_ESCURO = RGBColor(80, 90, 100)

# Status colors
VERDE_SUCESSO = RGBColor(40, 167, 69)
VERMELHO_ALERTA = RGBColor(220, 53, 69)
LARANJA_ATENCAO = RGBColor(255, 140, 0)


def add_premium_background(slide, style="navy"):
    """Adiciona fundo premium com gradiente simulado"""
    if style == "navy":
        # Camada base
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = AZUL_NAVY
        bg.line.fill.background()

        # Elemento decorativo sutil no canto
        accent = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2), Inches(6), Inches(6))
        accent.fill.solid()
        accent.fill.fore_color.rgb = AZUL_ROYAL
        accent.line.fill.background()

    elif style == "light":
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = CINZA_CLARO
        bg.line.fill.background()

    elif style == "gradient":
        # Simular gradiente com faixas
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = BRANCO
        bg.line.fill.background()

        # Faixa lateral elegante
        side = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.4), prs.slide_height)
        side.fill.solid()
        side.fill.fore_color.rgb = AZUL_NAVY
        side.line.fill.background()


def add_elegant_header(slide, title, subtitle="", with_line=True):
    """Adiciona header elegante e minimalista"""
    # Título principal
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title.upper()
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = AZUL_NAVY
    p.font.name = "Calibri Light"

    if subtitle:
        p = tf.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(14)
        p.font.color.rgb = CINZA_MEDIO
        p.space_before = Pt(5)

    if with_line:
        # Linha dourada embaixo do título
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.35), Inches(1.5), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = DOURADO
        line.line.fill.background()


def add_stat_card(slide, x, y, number, label, sublabel="", accent_color=None, width=2.5, height=1.8):
    """Card de estatística elegante com número grande"""
    if accent_color is None:
        accent_color = AZUL_ACCENT

    # Card com sombra simulada (retângulo atrás)
    shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x+0.05), Inches(y+0.05), Inches(width), Inches(height))
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = RGBColor(220, 220, 220)
    shadow.line.fill.background()

    # Card principal
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = BRANCO
    card.line.fill.background()

    # Barra de accent no topo
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(0.08))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent_color
    accent_bar.line.fill.background()

    # Número grande
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y + 0.25), Inches(width), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = accent_color
    p.alignment = PP_ALIGN.CENTER

    # Label
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y + 1.1), Inches(width), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = label.upper()
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = CINZA_ESCURO
    p.alignment = PP_ALIGN.CENTER

    if sublabel:
        txBox = slide.shapes.add_textbox(Inches(x), Inches(y + 1.4), Inches(width), Inches(0.3))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = sublabel
        p.font.size = Pt(9)
        p.font.color.rgb = CINZA_MEDIO
        p.alignment = PP_ALIGN.CENTER


def add_elegant_box(slide, x, y, width, height, title, items, icon_letter, color):
    """Box elegante com ícone circular"""
    # Container com borda sutil
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = BRANCO
    box.line.color.rgb = RGBColor(230, 230, 230)
    box.line.width = Pt(1)

    # Ícone circular no canto superior
    icon_size = 0.6
    icon = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.2), Inches(y + 0.2), Inches(icon_size), Inches(icon_size))
    icon.fill.solid()
    icon.fill.fore_color.rgb = color
    icon.line.fill.background()

    # Letra no ícone
    txBox = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.3), Inches(icon_size), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = icon_letter
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER

    # Título
    txBox = slide.shapes.add_textbox(Inches(x + 1), Inches(y + 0.25), Inches(width - 1.2), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = AZUL_NAVY

    # Itens
    txBox = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.9), Inches(width - 0.6), Inches(height - 1.1))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(12)
        p.font.color.rgb = CINZA_ESCURO
        p.space_after = Pt(4)


def add_process_step(slide, x, y, number, title, subtitle, is_active=False):
    """Step de processo elegante"""
    color = AZUL_ACCENT if is_active else CINZA_MEDIO

    # Círculo numerado
    size = 0.8
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()

    # Número
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y + 0.18), Inches(size), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER

    # Título abaixo
    txBox = slide.shapes.add_textbox(Inches(x - 0.3), Inches(y + 1), Inches(1.4), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title.upper()
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    txBox = slide.shapes.add_textbox(Inches(x - 0.5), Inches(y + 1.4), Inches(1.8), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(9)
    p.font.color.rgb = CINZA_ESCURO
    p.alignment = PP_ALIGN.CENTER


def add_connector_line(slide, x1, y, x2):
    """Linha conectora entre elementos"""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x1), Inches(y), Inches(x2-x1), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(200, 200, 200)
    line.line.fill.background()


def add_arrow_connector(slide, x, y, width=0.8):
    """Seta de conexão elegante"""
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(width), Inches(0.25))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = DOURADO
    arrow.line.fill.background()


# ============ CRIAR SLIDES ============

# ======== SLIDE 1 - CAPA ========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_premium_background(slide, "navy")

# Elemento decorativo geométrico
geo1 = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, Inches(9), Inches(-1), Inches(6), Inches(10))
geo1.fill.solid()
geo1.fill.fore_color.rgb = AZUL_ROYAL
geo1.line.fill.background()

# Linha dourada accent
gold_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.8), Inches(1.2), Inches(0.06))
gold_line.fill.solid()
gold_line.fill.fore_color.rgb = DOURADO
gold_line.line.fill.background()

# Título principal
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(3), Inches(8), Inches(1.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "PREVISÃO DE"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = BRANCO
p.font.name = "Calibri Light"

p = tf.add_paragraph()
p.text = "DEMANDA"
p.font.size = Pt(60)
p.font.bold = True
p.font.color.rgb = DOURADO
p.font.name = "Calibri Light"
p.space_before = Pt(0)

# Subtítulo
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(7), Inches(0.8))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Sistema Inteligente de Gestão de Pedidos"
p.font.size = Pt(18)
p.font.color.rgb = CINZA_CLARO

p = tf.add_paragraph()
p.text = "Proposta de Implantação Piloto"
p.font.size = Pt(14)
p.font.color.rgb = CINZA_MEDIO
p.space_before = Pt(10)

# Data
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(3), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Janeiro 2026"
p.font.size = Pt(12)
p.font.color.rgb = CINZA_MEDIO


# ======== SLIDE 2 - O PROBLEMA ========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_premium_background(slide, "gradient")
add_elegant_header(slide, "O Cenário Atual", "Desafios que impactam o resultado")

# Diagrama central - Problema no centro
center_x, center_y = 6.2, 3.8
center_size = 1.6

# Círculo central vermelho
center = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(center_x), Inches(center_y), Inches(center_size), Inches(center_size))
center.fill.solid()
center.fill.fore_color.rgb = VERMELHO_ALERTA
center.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(center_x), Inches(center_y + 0.4), Inches(center_size), Inches(0.8))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "GESTÃO"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "MANUAL"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

# Problemas ao redor
problemas = [
    (4.0, 2.0, "Rupturas", "Perda de vendas"),
    (8.0, 2.0, "Estoque Alto", "Capital parado"),
    (2.5, 4.2, "Sem Dados", "Decisão por intuição"),
    (9.5, 4.2, "Processo Lento", "Baixa produtividade"),
    (4.0, 6.0, "Pedidos Errados", "Custos extras"),
    (8.0, 6.0, "Sem Visibilidade", "Planejamento difícil"),
]

for x, y, title, desc in problemas:
    # Caixinha elegante
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2), Inches(0.9))
    box.fill.solid()
    box.fill.fore_color.rgb = BRANCO
    box.line.color.rgb = VERMELHO_ALERTA
    box.line.width = Pt(1.5)

    txBox = slide.shapes.add_textbox(Inches(x), Inches(y + 0.1), Inches(2), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = AZUL_NAVY
    p.alignment = PP_ALIGN.CENTER

    txBox = slide.shapes.add_textbox(Inches(x), Inches(y + 0.5), Inches(2), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(9)
    p.font.color.rgb = CINZA_MEDIO
    p.alignment = PP_ALIGN.CENTER


# ======== SLIDE 3 - A SOLUÇÃO (Fluxo Elegante) ========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_premium_background(slide, "gradient")
add_elegant_header(slide, "A Transformação", "De manual para inteligente")

# Três blocos principais: INPUT > PROCESSO > OUTPUT
# Bloco 1 - Input (Dados)
box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.2), Inches(3.2), Inches(4))
box1.fill.solid()
box1.fill.fore_color.rgb = BRANCO
box1.line.color.rgb = AZUL_ACCENT
box1.line.width = Pt(2)

# Header do box
header1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.2), Inches(3.2), Inches(0.7))
header1.fill.solid()
header1.fill.fore_color.rgb = AZUL_ACCENT
header1.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.35), Inches(3.2), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "DADOS"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

# Conteúdo box 1
dados_items = [
    ("10M+", "registros de venda"),
    ("3 anos", "de histórico"),
    ("701", "produtos ativos"),
    ("Todas", "as lojas")
]
y_pos = 3.1
for valor, label in dados_items:
    txBox = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(1.2), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = valor
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = AZUL_ACCENT

    txBox = slide.shapes.add_textbox(Inches(2.2), Inches(y_pos + 0.05), Inches(1.6), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(12)
    p.font.color.rgb = CINZA_ESCURO

    y_pos += 0.65

# Seta 1
add_arrow_connector(slide, 4.2, 4, 1)

# Bloco 2 - Processo (Sistema) - Central e maior
box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.4), Inches(1.8), Inches(3.5), Inches(4.8))
box2.fill.solid()
box2.fill.fore_color.rgb = AZUL_NAVY
box2.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(5.4), Inches(2.1), Inches(3.5), Inches(0.6))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "INTELIGENCIA"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = DOURADO
p.alignment = PP_ALIGN.CENTER

# Ícones/features do sistema
sistema_items = ["Algoritmos Estatísticos", "Machine Learning", "Análise de Sazonalidade", "Detecção de Tendências", "Otimização de Pedidos"]
y_pos = 2.8
for item in sistema_items:
    # Bolinha dourada
    bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.7), Inches(y_pos + 0.1), Inches(0.15), Inches(0.15))
    bullet.fill.solid()
    bullet.fill.fore_color.rgb = DOURADO
    bullet.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(6), Inches(y_pos), Inches(2.7), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = item
    p.font.size = Pt(13)
    p.font.color.rgb = BRANCO
    y_pos += 0.55

# Seta 2
add_arrow_connector(slide, 9.1, 4, 1)

# Bloco 3 - Output (Resultados)
box3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.3), Inches(2.2), Inches(2.5), Inches(4))
box3.fill.solid()
box3.fill.fore_color.rgb = BRANCO
box3.line.color.rgb = VERDE_SUCESSO
box3.line.width = Pt(2)

header3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.3), Inches(2.2), Inches(2.5), Inches(0.7))
header3.fill.solid()
header3.fill.fore_color.rgb = VERDE_SUCESSO
header3.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(10.3), Inches(2.35), Inches(2.5), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "RESULTADOS"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

resultados = ["Previsao 12 meses", "Pedidos otimizados", "Menos rupturas", "Mais vendas", "Equipe produtiva"]
y_pos = 3.1
for item in resultados:
    txBox = slide.shapes.add_textbox(Inches(10.4), Inches(y_pos), Inches(2.3), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"+ {item}"
    p.font.size = Pt(11)
    p.font.color.rgb = CINZA_ESCURO
    y_pos += 0.5


# ======== SLIDE 4 - MÓDULOS DO SISTEMA ========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_premium_background(slide, "gradient")
add_elegant_header(slide, "Funcionalidades", "Quatro pilares integrados")

# 4 módulos em grid 2x2
modulos = [
    (0.8, 1.8, "Previsao de Demanda", "P", AZUL_ACCENT,
     ["Forecast de 12 meses", "Analise por produto e loja", "Deteccao de sazonalidade", "Ajuste de tendencias"]),
    (6.9, 1.8, "Gestao de Pedidos", "G", VERDE_SUCESSO,
     ["Sugestao automatica", "Arredondamento por embalagem", "Calendario de entregas", "Tracking de pedidos"]),
    (0.8, 4.4, "Controle de Estoque", "E", LARANJA_ATENCAO,
     ["Posicao em tempo real", "Alertas de ruptura", "Balanceamento entre lojas", "Indicadores de giro"]),
    (6.9, 4.4, "Gestao de Fornecedores", "F", RGBColor(128, 90, 180),
     ["Nivel de servico (OTIF)", "Portfolio de produtos", "Historico de entregas", "Ranking de performance"])
]

for x, y, title, icon, color, items in modulos:
    add_elegant_box(slide, x, y, 6, 2.4, title, items, icon, color)


# ======== SLIDE 5 - ANTES vs DEPOIS ========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_premium_background(slide, "gradient")
add_elegant_header(slide, "A Mudanca", "Evolucao do processo")

# Lado ANTES
before_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5))
before_box.fill.solid()
before_box.fill.fore_color.rgb = RGBColor(255, 245, 245)
before_box.line.color.rgb = VERMELHO_ALERTA
before_box.line.width = Pt(2)

# Header ANTES
header_antes = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.5), Inches(0.8))
header_antes.fill.solid()
header_antes.fill.fore_color.rgb = VERMELHO_ALERTA
header_antes.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.95), Inches(5.5), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "HOJE"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

# Itens ANTES
antes_items = [
    "Calculo manual pela equipe",
    "Mesma formula para todos",
    "Decisao baseada em intuicao",
    "Sem analise de sazonalidade",
    "Processo reativo",
    "Sem visibilidade de demanda"
]
y_pos = 2.9
for item in antes_items:
    txBox = slide.shapes.add_textbox(Inches(1.1), Inches(y_pos), Inches(5), Inches(0.45))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"X   {item}"
    p.font.size = Pt(14)
    p.font.color.rgb = CINZA_ESCURO
    y_pos += 0.6

# Seta central grande
arrow_big = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.5), Inches(4), Inches(0.8), Inches(0.6))
arrow_big.fill.solid()
arrow_big.fill.fore_color.rgb = DOURADO
arrow_big.line.fill.background()

# Lado DEPOIS
after_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(1.8), Inches(5.5), Inches(5))
after_box.fill.solid()
after_box.fill.fore_color.rgb = RGBColor(240, 255, 240)
after_box.line.color.rgb = VERDE_SUCESSO
after_box.line.width = Pt(2)

# Header DEPOIS
header_depois = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.5), Inches(1.8), Inches(5.5), Inches(0.8))
header_depois.fill.solid()
header_depois.fill.fore_color.rgb = VERDE_SUCESSO
header_depois.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(7.5), Inches(1.95), Inches(5.5), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "COM O SISTEMA"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

# Itens DEPOIS
depois_items = [
    "Calculo automatico inteligente",
    "Personalizado por item e loja",
    "Decisao baseada em dados",
    "Sazonalidade calculada",
    "Gestao proativa",
    "Previsao de 12 meses"
]
y_pos = 2.9
for item in depois_items:
    txBox = slide.shapes.add_textbox(Inches(7.8), Inches(y_pos), Inches(5), Inches(0.45))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"+   {item}"
    p.font.size = Pt(14)
    p.font.color.rgb = CINZA_ESCURO
    y_pos += 0.6


# ======== SLIDE 6 - PROPOSTA DO PILOTO ========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_premium_background(slide, "gradient")
add_elegant_header(slide, "Proposta de Piloto", "Escopo controlado, resultados rapidos")

# KPIs principais
add_stat_card(slide, 0.8, 1.8, "4", "FORNECEDORES", "Selecionados", AZUL_ACCENT, 2.8, 1.6)
add_stat_card(slide, 4, 1.8, "701", "PRODUTOS", "No escopo", VERDE_SUCESSO, 2.8, 1.6)
add_stat_card(slide, 7.2, 1.8, "10M+", "REGISTROS", "Historico", LARANJA_ATENCAO, 2.8, 1.6)
add_stat_card(slide, 10.4, 1.8, "3 anos", "DADOS", "Disponiveis", RGBColor(128, 90, 180), 2.8, 1.6)

# Cards dos fornecedores
fornecedores = [
    ("FAME", "45 SKUs", "Ferragens", AZUL_ACCENT),
    ("ZAGONEL", "75 SKUs", "Chuveiros/Duchas", VERDE_SUCESSO),
    ("SILVANA", "311 SKUs", "Fechaduras", LARANJA_ATENCAO),
    ("STAM", "270 SKUs", "Fechaduras", RGBColor(128, 90, 180))
]

x_pos = 0.8
for nome, skus, cat, color in fornecedores:
    # Card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(4), Inches(3), Inches(2.2))
    card.fill.solid()
    card.fill.fore_color.rgb = BRANCO
    card.line.color.rgb = color
    card.line.width = Pt(2)

    # Barra colorida
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x_pos), Inches(4), Inches(0.12), Inches(2.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    # Nome
    txBox = slide.shapes.add_textbox(Inches(x_pos + 0.3), Inches(4.2), Inches(2.5), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = nome
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = AZUL_NAVY

    # SKUs
    txBox = slide.shapes.add_textbox(Inches(x_pos + 0.3), Inches(4.8), Inches(2.5), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = skus
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = color

    # Categoria
    txBox = slide.shapes.add_textbox(Inches(x_pos + 0.3), Inches(5.4), Inches(2.5), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = cat
    p.font.size = Pt(12)
    p.font.color.rgb = CINZA_MEDIO

    x_pos += 3.2

# Rodapé
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(12), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Todas as lojas da rede   |   Sistema pronto para uso   |   Dados ja carregados e validados"
p.font.size = Pt(12)
p.font.color.rgb = CINZA_MEDIO
p.alignment = PP_ALIGN.CENTER


# ======== SLIDE 7 - BENEFÍCIOS ========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_premium_background(slide, "gradient")
add_elegant_header(slide, "Beneficios Esperados", "Impacto mensuravel no resultado")

# Grid de benefícios com números grandes
beneficios = [
    ("-30%", "REDUCAO DE RUPTURA", "Mais produtos disponiveis, mais vendas", VERDE_SUCESSO),
    ("+40%", "ACURACIDADE", "Pedidos mais precisos e assertivos", AZUL_ACCENT),
    ("+25%", "EFICIENCIA DE ESTOQUE", "Melhor balanceamento entre lojas", LARANJA_ATENCAO),
    ("+50%", "PRODUTIVIDADE", "Menos tempo em calculos manuais", RGBColor(128, 90, 180))
]

x_pos = 0.8
for valor, titulo, desc, color in beneficios:
    # Card grande
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(1.8), Inches(3), Inches(3.8))
    card.fill.solid()
    card.fill.fore_color.rgb = BRANCO
    card.line.color.rgb = RGBColor(230, 230, 230)
    card.line.width = Pt(1)

    # Círculo com número
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x_pos + 0.7), Inches(2.1), Inches(1.6), Inches(1.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(x_pos + 0.7), Inches(2.5), Inches(1.6), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = valor
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER

    # Título
    txBox = slide.shapes.add_textbox(Inches(x_pos), Inches(3.9), Inches(3), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = AZUL_NAVY
    p.alignment = PP_ALIGN.CENTER

    # Descrição
    txBox = slide.shapes.add_textbox(Inches(x_pos + 0.2), Inches(4.4), Inches(2.6), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(10)
    p.font.color.rgb = CINZA_MEDIO
    p.alignment = PP_ALIGN.CENTER

    x_pos += 3.15

# Box de impacto financeiro
impact_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6), Inches(12.5), Inches(0.9))
impact_box.fill.solid()
impact_box.fill.fore_color.rgb = AZUL_NAVY
impact_box.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(12.5), Inches(0.6))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "IMPACTO FINANCEIRO:   Reducao de custos logisticos   |   Otimizacao do capital de giro   |   Melhor fluxo de caixa"
p.font.size = Pt(14)
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER


# ======== SLIDE 8 - ROADMAP ========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_premium_background(slide, "gradient")
add_elegant_header(slide, "Roadmap de Expansao", "Crescimento gradual e controlado")

# Linha de conexão principal
add_connector_line(slide, 1.6, 3.35, 11.8)

# Fases do roadmap
fases = [
    (1.2, "1", "PILOTO", "4 fornecedores\n701 SKUs", True),
    (4.2, "2", "EXPANSAO I", "+8 fornecedores\n~2.000 SKUs", False),
    (7.2, "3", "EXPANSAO II", "+11 fornecedores\n~5.000 SKUs", False),
    (10.2, "4", "COMPLETO", "23 categorias\nTodos produtos", False)
]

for x, num, titulo, desc, active in fases:
    add_process_step(slide, x, 2.9, num, titulo, desc, active)

# Caixa de critérios
criteria_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(5.2), Inches(9.3), Inches(1.5))
criteria_box.fill.solid()
criteria_box.fill.fore_color.rgb = BRANCO
criteria_box.line.color.rgb = RGBColor(230, 230, 230)
criteria_box.line.width = Pt(1)

txBox = slide.shapes.add_textbox(Inches(2), Inches(5.4), Inches(9.3), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "CRITERIOS DE EXPANSAO"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = AZUL_NAVY
p.alignment = PP_ALIGN.CENTER

criterios = ["Resultados positivos na fase anterior", "Dados historicos disponiveis", "Capacidade operacional da equipe"]
txBox = slide.shapes.add_textbox(Inches(2.5), Inches(5.9), Inches(8.3), Inches(0.8))
tf = txBox.text_frame
for i, crit in enumerate(criterios):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = f"     {crit}"
    p.font.size = Pt(11)
    p.font.color.rgb = CINZA_ESCURO
    p.alignment = PP_ALIGN.CENTER


# ======== SLIDE 9 - CRONOGRAMA ========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_premium_background(slide, "gradient")
add_elegant_header(slide, "Cronograma do Piloto", "12 semanas de implantacao")

# Gantt chart elegante
fases_gantt = [
    ("Setup e Configuracao", "Sem 1-2", 0, 2, AZUL_ACCENT),
    ("Treinamento da Equipe", "Sem 3-4", 2, 2, VERDE_SUCESSO),
    ("Operacao Assistida", "Sem 5-8", 4, 4, LARANJA_ATENCAO),
    ("Operacao Autonoma", "Sem 9-12", 8, 4, RGBColor(128, 90, 180))
]

# Grid de semanas
for i in range(13):
    x = 4.5 + i * 0.65
    # Linha vertical sutil
    vline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(2), Inches(0.01), Inches(3.5))
    vline.fill.solid()
    vline.fill.fore_color.rgb = RGBColor(230, 230, 230)
    vline.line.fill.background()

    # Número da semana
    if i > 0:
        txBox = slide.shapes.add_textbox(Inches(x - 0.15), Inches(1.7), Inches(0.4), Inches(0.3))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = str(i)
        p.font.size = Pt(9)
        p.font.color.rgb = CINZA_MEDIO
        p.alignment = PP_ALIGN.CENTER

y_pos = 2.2
for fase, periodo, start, duration, color in fases_gantt:
    # Label da fase
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos + 0.1), Inches(3.5), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = fase
    p.font.size = Pt(12)
    p.font.color.rgb = AZUL_NAVY

    # Barra do Gantt
    bar_x = 4.5 + start * 0.65
    bar_width = duration * 0.65

    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(bar_x), Inches(y_pos), Inches(bar_width), Inches(0.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    y_pos += 0.8

# Marcos importantes
marcos_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.5), Inches(12.5), Inches(1.3))
marcos_box.fill.solid()
marcos_box.fill.fore_color.rgb = RGBColor(250, 250, 250)
marcos_box.line.color.rgb = DOURADO
marcos_box.line.width = Pt(2)

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(5.7), Inches(12.5), Inches(0.3))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "MARCOS DO PROJETO"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = DOURADO
p.alignment = PP_ALIGN.CENTER

marcos = [
    ("Semana 4", "Primeira previsao gerada"),
    ("Semana 8", "Primeiro ciclo completo de pedidos"),
    ("Semana 12", "Avaliacao de resultados e decisao")
]

txBox = slide.shapes.add_textbox(Inches(1.5), Inches(6.1), Inches(10.5), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "   |   ".join([f"{m[0]}: {m[1]}" for m in marcos])
p.font.size = Pt(11)
p.font.color.rgb = CINZA_ESCURO
p.alignment = PP_ALIGN.CENTER


# ======== SLIDE 10 - POR QUE AGORA ========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_premium_background(slide, "gradient")
add_elegant_header(slide, "Por Que Agora?", "O momento e a oportunidade")

# 5 razões em círculos conectados
razoes = [
    ("DADOS\nPRONTOS", "Historico carregado\ne validado", AZUL_ACCENT),
    ("SISTEMA\nFUNCIONAL", "Pronto para\nuso imediato", VERDE_SUCESSO),
    ("BAIXO\nINVESTIMENTO", "Infraestrutura\nja disponivel", LARANJA_ATENCAO),
    ("ALTO\nRETORNO", "Potencial de\nganho comprovado", RGBColor(128, 90, 180)),
    ("RISCO\nCONTROLADO", "Piloto permite\najustes rapidos", RGBColor(220, 100, 100))
]

x_pos = 0.5
for titulo, desc, color in razoes:
    # Círculo grande
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x_pos + 0.4), Inches(2), Inches(1.7), Inches(1.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()

    # Texto no círculo
    txBox = slide.shapes.add_textbox(Inches(x_pos + 0.4), Inches(2.4), Inches(1.7), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER

    # Descrição abaixo
    txBox = slide.shapes.add_textbox(Inches(x_pos), Inches(3.9), Inches(2.5), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(10)
    p.font.color.rgb = CINZA_ESCURO
    p.alignment = PP_ALIGN.CENTER

    x_pos += 2.55

# Call to Action grande
cta_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(5.3), Inches(8.3), Inches(1.5))
cta_box.fill.solid()
cta_box.fill.fore_color.rgb = AZUL_NAVY
cta_box.line.fill.background()

# Linha dourada acima do CTA
gold_accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(5.2), Inches(2.3), Inches(0.06))
gold_accent.fill.solid()
gold_accent.fill.fore_color.rgb = DOURADO
gold_accent.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(2.5), Inches(5.5), Inches(8.3), Inches(0.7))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "VAMOS COMECAR?"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

txBox = slide.shapes.add_textbox(Inches(2.5), Inches(6.2), Inches(8.3), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "O piloto pode iniciar imediatamente apos aprovacao"
p.font.size = Pt(14)
p.font.color.rgb = CINZA_CLARO
p.alignment = PP_ALIGN.CENTER


# ======== SLIDE 11 - ENCERRAMENTO ========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_premium_background(slide, "navy")

# Elemento geométrico
geo = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, Inches(-2), Inches(0), Inches(8), Inches(10))
geo.fill.solid()
geo.fill.fore_color.rgb = AZUL_ROYAL
geo.line.fill.background()

# Linha dourada decorativa
gold_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.8), Inches(3.2), Inches(1.8), Inches(0.06))
gold_line.fill.solid()
gold_line.fill.fore_color.rgb = DOURADO
gold_line.line.fill.background()

# Obrigado
txBox = slide.shapes.add_textbox(Inches(0), Inches(3.4), Inches(13.333), Inches(1.2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "OBRIGADO"
p.font.size = Pt(64)
p.font.bold = True
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

# Subtítulo
txBox = slide.shapes.add_textbox(Inches(0), Inches(4.8), Inches(13.333), Inches(0.6))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Duvidas e proximos passos"
p.font.size = Pt(20)
p.font.color.rgb = CINZA_CLARO
p.alignment = PP_ALIGN.CENTER


# ============ SALVAR ============
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'Apresentacao_Previsao_Demanda_v3.pptx')
prs.save(output_path)
print(f'Apresentacao Premium salva em: {output_path}')
print('Versao 3 - Design de Consultoria')
