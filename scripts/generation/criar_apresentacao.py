# -*- coding: utf-8 -*-
"""Script para criar apresentação PowerPoint"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Criar apresentação
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Cores
AZUL_ESCURO = RGBColor(0, 51, 102)
AZUL_CLARO = RGBColor(0, 112, 192)
VERDE = RGBColor(0, 128, 0)
LARANJA = RGBColor(255, 153, 0)

def add_title_slide(title, subtitle=""):
    """Adiciona slide de título"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Fundo azul
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = AZUL_ESCURO
    shape.line.fill.background()

    # Título
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        p = tf.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(title, bullets, subtitle=""):
    """Adiciona slide com conteúdo"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Barra de título
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = AZUL_ESCURO
    shape.line.fill.background()

    # Título
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Subtítulo se houver
    y_start = Inches(1.5)
    if subtitle:
        txBox = slide.shapes.add_textbox(Inches(0.5), y_start, Inches(12.333), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = AZUL_CLARO
        y_start = Inches(2.2)

    # Bullets
    if bullets:
        txBox = slide.shapes.add_textbox(Inches(0.5), y_start, Inches(12.333), Inches(5))
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(51, 51, 51)
            p.space_after = Pt(12)
            p.level = 0

    return slide

def add_two_column_slide(title, left_title, left_bullets, right_title, right_bullets):
    """Adiciona slide com duas colunas"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Barra de título
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = AZUL_ESCURO
    shape.line.fill.background()

    # Título
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Coluna esquerda - título
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = AZUL_CLARO

    # Coluna esquerda - bullets
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(6), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(left_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(8)

    # Coluna direita - título
    txBox = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(6), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = AZUL_CLARO

    # Coluna direita - bullets
    txBox = slide.shapes.add_textbox(Inches(6.8), Inches(2.1), Inches(6), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(right_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(8)

    return slide

def add_table_slide(title, headers, rows):
    """Adiciona slide com tabela"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Barra de título
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = AZUL_ESCURO
    shape.line.fill.background()

    # Título
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Tabela
    num_cols = len(headers)
    num_rows = len(rows) + 1

    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.6), Inches(12.333), Inches(0.5 * num_rows)).table

    # Headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = AZUL_ESCURO
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    # Rows
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.alignment = PP_ALIGN.CENTER
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)

    return slide

# ============ CRIAR SLIDES ============

# Slide 1 - Capa
add_title_slide(
    "SISTEMA DE PREVISÃO DE DEMANDA\nE GESTÃO INTELIGENTE DE PEDIDOS",
    "Proposta de Implantação Piloto | Janeiro/2026"
)

# Slide 2 - Agenda
add_content_slide("AGENDA", [
    "1. Contexto Atual",
    "2. A Solução Proposta",
    "3. Funcionalidades do Sistema",
    "4. Proposta de Implantação Piloto",
    "5. Benefícios Esperados",
    "6. Roadmap de Expansão",
    "7. Próximos Passos"
])

# Slide 3 - Contexto Atual
add_content_slide("CONTEXTO ATUAL", [
    "• Cálculo manual de pedidos pela equipe comercial",
    "• Método simplificado e padronizado para todos os itens e lojas",
    "• Sem previsão de demanda estruturada",
    "• Decisões baseadas em experiência individual",
    "• Mesma regra aplicada independente de:",
    "    - Sazonalidade do produto",
    "    - Comportamento de cada loja",
    "    - Histórico de vendas específico"
], "Como funciona hoje:")

# Slide 4 - Desafios Atuais
add_table_slide("DESAFIOS ATUAIS",
    ["Desafio", "Impacto"],
    [
        ["Rupturas frequentes", "Perda de vendas e clientes"],
        ["Excesso de estoque", "Capital de giro imobilizado"],
        ["Pedidos imprecisos", "Custos logísticos elevados"],
        ["Gestão manual", "Baixa produtividade da equipe"],
        ["Sem visibilidade", "Dificuldade de planejamento"],
        ["Fornecedores", "Falta de controle de nível de serviço"]
    ]
)

# Slide 5 - A Solução
add_content_slide("A SOLUÇÃO", [
    "Plataforma inteligente que utiliza:",
    "",
    "• Algoritmos estatísticos avançados",
    "• Histórico de vendas de 3+ anos",
    "• Sazonalidade por produto/loja",
    "• Machine Learning para refinamento",
    "",
    "Resultado: Previsões precisas e pedidos otimizados"
], "Sistema de Previsão de Demanda")

# Slide 6 - Funcionalidades
add_two_column_slide("FUNCIONALIDADES PRINCIPAIS",
    "Previsão de Demanda",
    [
        "• Previsão diária/semanal/mensal",
        "• Por produto, loja e fornecedor",
        "• Considerando sazonalidade",
        "• Análise de tendências"
    ],
    "Gestão de Pedidos",
    [
        "• Pedidos planejados automáticos",
        "• Pedidos de reabastecimento",
        "• Arredondamento por embalagem",
        "• Calendário de entregas"
    ]
)

# Slide 7 - Módulos do Sistema
add_two_column_slide("MÓDULOS DO SISTEMA",
    "Gestão de Estoque",
    [
        "• Posição atual por loja",
        "• Alertas de ruptura",
        "• Balanceamento entre lojas",
        "• Indicadores de giro"
    ],
    "Gestão de Fornecedores",
    [
        "• Portfólio de produtos",
        "• Nível de serviço (OTIF)",
        "• Tracking de pedidos",
        "• Performance de entrega"
    ]
)

# Slide 8 - Proposta de Implantação
add_content_slide("PROPOSTA DE IMPLANTAÇÃO", [
    "Estratégia: PILOTO CONTROLADO",
    "",
    "Fase 1 - Piloto (3 meses)",
    "",
    "• 4 fornecedores selecionados:",
    "    - FAME (Ferragens)",
    "    - ZAGONEL (Chuveiros/Duchas)",
    "    - SILVANA (Fechaduras)",
    "    - STAM (Fechaduras)",
    "",
    "• 701 produtos | Todas as lojas da rede"
])

# Slide 9 - Dados do Piloto
add_table_slide("DADOS DO PILOTO - Base Disponível",
    ["Fornecedor", "Produtos", "Registros de Venda", "Período"],
    [
        ["FAME", "45", "804.240", "2023-2026"],
        ["ZAGONEL", "75", "1.324.762", "2023-2026"],
        ["SILVANA", "311", "3.000.000+", "2023-2026"],
        ["STAM", "270", "3.000.000+", "2023-2026"],
        ["TOTAL", "701", "10+ milhões", "3 anos"]
    ]
)

# Slide 10 - Escopo do Piloto
add_two_column_slide("ESCOPO DO PILOTO",
    "Cálculo de Demanda",
    [
        "• Previsão 12 meses (2026)",
        "• Análise de sazonalidade",
        "• Tendências por loja",
        "• Ajustes por eventos"
    ],
    "Gestão Operacional",
    [
        "• Sugestão automática de pedidos",
        "• Monitoramento de ruptura",
        "• Controle de entregas",
        "• Nível de serviço fornecedores"
    ]
)

# Slide 11 - Benefícios Quantitativos
add_table_slide("BENEFÍCIOS ESPERADOS - Quantitativos",
    ["Benefício", "Meta Piloto", "Impacto Financeiro"],
    [
        ["Redução de Ruptura", "-30%", "Aumento de Vendas"],
        ["Acuracidade Pedidos", "+40%", "Redução Custos Logísticos"],
        ["Balanceamento Estoque", "+25%", "Redução Capital de Giro"],
        ["Produtividade Equipe", "+50%", "Redução Horas Trabalhadas"],
        ["Nível Serviço Fornec.", "+20%", "Maior Disponibilidade"]
    ]
)

# Slide 12 - Benefícios Qualitativos
add_two_column_slide("BENEFÍCIOS ESPERADOS - Qualitativos",
    "Para a Operação",
    [
        "• Decisões baseadas em dados",
        "• Padronização de processos",
        "• Visibilidade em tempo real",
        "• Antecipação de problemas"
    ],
    "Para o Financeiro",
    [
        "• Melhor capital de giro",
        "• Otimização fluxo de caixa",
        "• Redução custos logísticos",
        "• Menor imobilização estoque"
    ]
)

# Slide 13 - Roadmap
add_content_slide("ROADMAP DE EXPANSÃO", [
    "FASE 1 (3 meses) → FASE 2 (3 meses) → FASE 3 (6 meses) → FASE 4",
    "",
    "• Fase 1: PILOTO - 4 fornecedores, 701 SKUs",
    "• Fase 2: EXPANSÃO - +8 fornecedores, ~2.000 SKUs",
    "• Fase 3: EXPANSÃO - +11 fornecedores, ~5.000 SKUs",
    "• Fase 4: TODAS - 23 categorias completas",
    "",
    "Critérios de Expansão:",
    "• Resultados positivos na fase anterior",
    "• Dados históricos disponíveis",
    "• Capacidade da equipe"
])

# Slide 14 - Recursos
add_two_column_slide("INVESTIMENTO E RECURSOS",
    "Infraestrutura (já disponível)",
    [
        "• Servidor de banco de dados",
        "• Sistema web desenvolvido",
        "• Dados históricos carregados",
        "• Ambiente de produção pronto"
    ],
    "Equipe Necessária",
    [
        "• 1 Analista de Dados (parcial)",
        "• Equipe Comercial (usuários)",
        "• TI (suporte técnico)",
        "• Gestão (acompanhamento)"
    ]
)

# Slide 15 - Cronograma
add_content_slide("CRONOGRAMA DO PILOTO", [
    "Semana 1-2: Setup e Configuração",
    "Semana 3-4: Treinamento da Equipe",
    "Semana 5-8: Operação Assistida",
    "Semana 9-12: Operação Autônoma + Avaliação",
    "",
    "Marcos Importantes:",
    "• Marco 1 (Sem 4): Primeira previsão gerada",
    "• Marco 2 (Sem 8): Primeiro ciclo completo de pedidos",
    "• Marco 3 (Sem 12): Avaliação de resultados"
])

# Slide 16 - KPIs
add_content_slide("INDICADORES DE SUCESSO", [
    "KPIs do Piloto:",
    "",
    "1. Acuracidade da Previsão → Meta: >80% (MAPE < 20%)",
    "2. Redução de Ruptura → Meta: -30% vs. período anterior",
    "3. Nível de Serviço → Meta: >95% de disponibilidade",
    "4. Adoção pela Equipe → Meta: 100% dos pedidos via sistema",
    "5. Satisfação do Usuário → Meta: NPS > 7"
])

# Slide 17 - Caso de Sucesso
add_content_slide("CASO DE SUCESSO DEMONSTRADO", [
    "Base carregada e validada:",
    "",
    "• 3 anos de histórico (2023-2025)",
    "• 10+ milhões de registros de vendas",
    "• 701 produtos dos 4 fornecedores",
    "• Dados de estoque atualizados",
    "• Cadastro de produtos completo",
    "",
    "Sistema funcional:",
    "• Previsão de demanda operacional",
    "• Dashboards implementados",
    "• Relatórios disponíveis"
])

# Slide 18 - Próximos Passos
add_two_column_slide("PRÓXIMOS PASSOS",
    "Ações Imediatas",
    [
        "1. Aprovação desta proposta",
        "2. Definição da equipe piloto",
        "3. Treinamento inicial",
        "4. Início da operação piloto"
    ],
    "Entregáveis do Piloto",
    [
        "• Previsão de demanda 12 meses",
        "• Sugestões de pedidos semanais",
        "• Relatórios de acompanhamento",
        "• Avaliação de resultados"
    ]
)

# Slide 19 - Conclusão
add_content_slide("POR QUE AGORA?", [
    "• Dados disponíveis e validados",
    "• Sistema pronto para uso",
    "• Baixo investimento inicial",
    "• Alto potencial de retorno",
    "• Risco controlado via piloto",
    "",
    "",
    "O momento é agora!"
])

# Slide 20 - Final
add_title_slide("OBRIGADO!", "Perguntas?")

# Salvar
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'Apresentacao_Previsao_Demanda.pptx')
prs.save(output_path)
print(f'Apresentação salva em: {output_path}')
