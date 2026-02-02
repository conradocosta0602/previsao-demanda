# -*- coding: utf-8 -*-
"""
Script para criar apresentação PowerPoint para engajamento de fornecedores
Sistema de Previsão de Demanda e Gestão de Abastecimento - Ferreira Costa
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
import os

# Criar apresentação widescreen
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# =============================================================================
# PALETA DE CORES FERREIRA COSTA
# =============================================================================
AZUL_NAVY = RGBColor(15, 32, 65)        # Fundo principal
AZUL_ACCENT = RGBColor(0, 120, 200)     # Destaque principal
VERDE_SUCESSO = RGBColor(40, 167, 69)   # Sucesso/Positivo
LARANJA_ATENCAO = RGBColor(255, 140, 0) # Atenção/Alerta
ROXO = RGBColor(128, 90, 180)           # Diferenciação
VERMELHO = RGBColor(220, 53, 69)        # Crítico
BRANCO = RGBColor(255, 255, 255)
CINZA_CLARO = RGBColor(245, 247, 250)
CINZA_TEXTO = RGBColor(100, 110, 120)
CINZA_MEDIO = RGBColor(160, 170, 180)

# =============================================================================
# FUNÇÕES AUXILIARES
# =============================================================================

def add_header_bar(slide, title, subtitle=None):
    """Adiciona barra de header padronizada"""
    # Barra azul navy no topo
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = AZUL_NAVY
    header.line.fill.background()

    # Título
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BRANCO

    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.75), Inches(12), Inches(0.3))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = CINZA_MEDIO

def add_card(slide, x, y, width, height, title, items, color, icon_text=None):
    """Adiciona card com borda colorida"""
    # Fundo do card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = BRANCO
    card.line.color.rgb = color
    card.line.width = Pt(2)

    # Barra superior colorida
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    # Ícone se fornecido
    y_offset = 0.25
    if icon_text:
        icon = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.15), Inches(y + 0.25), Inches(0.5), Inches(0.5))
        icon.fill.solid()
        icon.fill.fore_color.rgb = color
        icon.line.fill.background()

        txIcon = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.32), Inches(0.5), Inches(0.4))
        tfIcon = txIcon.text_frame
        pIcon = tfIcon.paragraphs[0]
        pIcon.text = icon_text
        pIcon.font.size = Pt(16)
        pIcon.font.bold = True
        pIcon.font.color.rgb = BRANCO
        pIcon.alignment = PP_ALIGN.CENTER
        y_offset = 0.85

    # Título do card
    txBox = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + y_offset), Inches(width - 0.3), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = AZUL_NAVY

    # Items
    txItems = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + y_offset + 0.35), Inches(width - 0.3), Inches(height - y_offset - 0.5))
    tfItems = txItems.text_frame
    tfItems.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tfItems.paragraphs[0]
        else:
            p = tfItems.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(11)
        p.font.color.rgb = CINZA_TEXTO
        p.space_after = Pt(4)

def add_metric_box(slide, x, y, value, label, color, width=2.0):
    """Adiciona box de métrica destacada"""
    # Fundo
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(1.3))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()

    # Valor
    txVal = slide.shapes.add_textbox(Inches(x), Inches(y + 0.15), Inches(width), Inches(0.6))
    tfVal = txVal.text_frame
    pVal = tfVal.paragraphs[0]
    pVal.text = value
    pVal.font.size = Pt(32)
    pVal.font.bold = True
    pVal.font.color.rgb = BRANCO
    pVal.alignment = PP_ALIGN.CENTER

    # Label
    txLbl = slide.shapes.add_textbox(Inches(x), Inches(y + 0.75), Inches(width), Inches(0.5))
    tfLbl = txLbl.text_frame
    tfLbl.word_wrap = True
    pLbl = tfLbl.paragraphs[0]
    pLbl.text = label
    pLbl.font.size = Pt(11)
    pLbl.font.color.rgb = BRANCO
    pLbl.alignment = PP_ALIGN.CENTER

def add_flow_arrow(slide, x1, y, x2, color=AZUL_ACCENT):
    """Adiciona seta de fluxo"""
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x1), Inches(y), Inches(x2 - x1), Inches(0.3))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()

def add_process_step(slide, x, y, number, title, description, color):
    """Adiciona passo de processo numerado"""
    # Círculo numerado
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()

    txNum = slide.shapes.add_textbox(Inches(x), Inches(y + 0.12), Inches(0.6), Inches(0.4))
    tfNum = txNum.text_frame
    pNum = tfNum.paragraphs[0]
    pNum.text = str(number)
    pNum.font.size = Pt(20)
    pNum.font.bold = True
    pNum.font.color.rgb = BRANCO
    pNum.alignment = PP_ALIGN.CENTER

    # Título
    txTitle = slide.shapes.add_textbox(Inches(x - 0.3), Inches(y + 0.7), Inches(1.2), Inches(0.4))
    tfTitle = txTitle.text_frame
    tfTitle.word_wrap = True
    pTitle = tfTitle.paragraphs[0]
    pTitle.text = title
    pTitle.font.size = Pt(11)
    pTitle.font.bold = True
    pTitle.font.color.rgb = AZUL_NAVY
    pTitle.alignment = PP_ALIGN.CENTER

    # Descrição
    txDesc = slide.shapes.add_textbox(Inches(x - 0.4), Inches(y + 1.05), Inches(1.4), Inches(0.8))
    tfDesc = txDesc.text_frame
    tfDesc.word_wrap = True
    pDesc = tfDesc.paragraphs[0]
    pDesc.text = description
    pDesc.font.size = Pt(9)
    pDesc.font.color.rgb = CINZA_TEXTO
    pDesc.alignment = PP_ALIGN.CENTER

# =============================================================================
# SLIDE 1 - CAPA
# =============================================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])

# Fundo gradiente (simulado com retângulo)
bg = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = AZUL_NAVY
bg.line.fill.background()

# Linha decorativa
line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3.2), Inches(12.333), Inches(0.02))
line.fill.solid()
line.fill.fore_color.rgb = AZUL_ACCENT
line.line.fill.background()

# Título principal
txTitle = slide1.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12.333), Inches(1.0))
tf = txTitle.text_frame
p = tf.paragraphs[0]
p.text = "PARCERIA PARA EXCELÊNCIA"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

# Subtítulo
txSub = slide1.shapes.add_textbox(Inches(0.5), Inches(3.4), Inches(12.333), Inches(0.8))
tf2 = txSub.text_frame
p2 = tf2.paragraphs[0]
p2.text = "Sistema de Gestão Inteligente de Abastecimento"
p2.font.size = Pt(28)
p2.font.color.rgb = AZUL_ACCENT
p2.alignment = PP_ALIGN.CENTER

# Linha inferior
txLine = slide1.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(0.5))
tf3 = txLine.text_frame
p3 = tf3.paragraphs[0]
p3.text = "Elevando o Nível de Serviço Juntos"
p3.font.size = Pt(18)
p3.font.color.rgb = CINZA_MEDIO
p3.alignment = PP_ALIGN.CENTER

# Logo placeholder
txLogo = slide1.shapes.add_textbox(Inches(5.5), Inches(6.0), Inches(2.333), Inches(0.5))
tf4 = txLogo.text_frame
p4 = tf4.paragraphs[0]
p4.text = "FERREIRA COSTA"
p4.font.size = Pt(16)
p4.font.bold = True
p4.font.color.rgb = BRANCO
p4.alignment = PP_ALIGN.CENTER

# =============================================================================
# SLIDE 2 - NOSSO OBJETIVO COMUM
# =============================================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])

# Fundo
bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg2.fill.solid()
bg2.fill.fore_color.rgb = CINZA_CLARO
bg2.line.fill.background()

add_header_bar(slide2, "NOSSO OBJETIVO COMUM", "Melhorar a qualidade do pedido e elevar o nível de serviço")

# Diagrama central - Objetivo
central_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(2.0), Inches(4.333), Inches(1.2))
central_box.fill.solid()
central_box.fill.fore_color.rgb = AZUL_NAVY
central_box.line.fill.background()

txCentral = slide2.shapes.add_textbox(Inches(4.5), Inches(2.3), Inches(4.333), Inches(0.8))
tfCentral = txCentral.text_frame
pCentral = tfCentral.paragraphs[0]
pCentral.text = "PEDIDO DE QUALIDADE"
pCentral.font.size = Pt(22)
pCentral.font.bold = True
pCentral.font.color.rgb = BRANCO
pCentral.alignment = PP_ALIGN.CENTER

# Benefícios do fornecedor (esquerda)
add_card(slide2, 0.5, 3.5, 3.8, 2.8, "BENEFÍCIOS PARA O FORNECEDOR",
         ["Previsibilidade de demanda", "Redução de devoluções", "Planejamento de produção", "Relacionamento de longo prazo"],
         VERDE_SUCESSO, "F")

# Benefícios Ferreira Costa (centro)
add_card(slide2, 4.75, 3.5, 3.8, 2.8, "BENEFÍCIOS PARA FERREIRA COSTA",
         ["Estoque balanceado", "Menos rupturas nas gôndolas", "Satisfação do cliente final", "Eficiência operacional"],
         AZUL_ACCENT, "FC")

# Benefícios cliente final (direita)
add_card(slide2, 9.0, 3.5, 3.8, 2.8, "BENEFÍCIOS PARA O CLIENTE",
         ["Produto disponível", "Preço competitivo", "Variedade mantida", "Experiência de compra"],
         LARANJA_ATENCAO, "C")

# Setas conectoras
arrow1 = slide2.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(2.0), Inches(2.5), Inches(2.3), Inches(0.2))
arrow1.fill.solid()
arrow1.fill.fore_color.rgb = VERDE_SUCESSO
arrow1.line.fill.background()

arrow2 = slide2.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(9.0), Inches(2.5), Inches(2.3), Inches(0.2))
arrow2.fill.solid()
arrow2.fill.fore_color.rgb = LARANJA_ATENCAO
arrow2.line.fill.background()
arrow2.rotation = 180

# =============================================================================
# SLIDE 3 - O QUE É UM PEDIDO DE QUALIDADE
# =============================================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])

bg3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg3.fill.solid()
bg3.fill.fore_color.rgb = CINZA_CLARO
bg3.line.fill.background()

add_header_bar(slide3, "O QUE É UM PEDIDO DE QUALIDADE?", "Características que definem um pedido bem elaborado")

# 4 pilares
pilares = [
    ("QUANTIDADE CERTA", "Baseada em previsão estatística com 6 algoritmos", AZUL_ACCENT, "Q"),
    ("MOMENTO CERTO", "Ciclo fixo com lead time respeitado", VERDE_SUCESSO, "T"),
    ("PRODUTO CERTO", "Mix adequado por curva ABC", LARANJA_ATENCAO, "P"),
    ("ENTREGA INTEGRAL", "Compromisso de 100% do pedido", ROXO, "E")
]

x_positions = [0.5, 3.6, 6.7, 9.8]
for i, (titulo, desc, cor, icone) in enumerate(pilares):
    # Box do pilar
    box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_positions[i]), Inches(1.8), Inches(2.8), Inches(2.5))
    box.fill.solid()
    box.fill.fore_color.rgb = BRANCO
    box.line.color.rgb = cor
    box.line.width = Pt(3)

    # Ícone circular
    circle = slide3.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x_positions[i] + 0.9), Inches(2.0), Inches(1.0), Inches(1.0))
    circle.fill.solid()
    circle.fill.fore_color.rgb = cor
    circle.line.fill.background()

    txIcon = slide3.shapes.add_textbox(Inches(x_positions[i] + 0.9), Inches(2.25), Inches(1.0), Inches(0.6))
    tfIcon = txIcon.text_frame
    pIcon = tfIcon.paragraphs[0]
    pIcon.text = icone
    pIcon.font.size = Pt(28)
    pIcon.font.bold = True
    pIcon.font.color.rgb = BRANCO
    pIcon.alignment = PP_ALIGN.CENTER

    # Título
    txT = slide3.shapes.add_textbox(Inches(x_positions[i] + 0.1), Inches(3.1), Inches(2.6), Inches(0.5))
    tfT = txT.text_frame
    tfT.word_wrap = True
    pT = tfT.paragraphs[0]
    pT.text = titulo
    pT.font.size = Pt(13)
    pT.font.bold = True
    pT.font.color.rgb = AZUL_NAVY
    pT.alignment = PP_ALIGN.CENTER

    # Descrição
    txD = slide3.shapes.add_textbox(Inches(x_positions[i] + 0.1), Inches(3.55), Inches(2.6), Inches(0.8))
    tfD = txD.text_frame
    tfD.word_wrap = True
    pD = tfD.paragraphs[0]
    pD.text = desc
    pD.font.size = Pt(10)
    pD.font.color.rgb = CINZA_TEXTO
    pD.alignment = PP_ALIGN.CENTER

# Box de resultado
result_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(4.8), Inches(8.333), Inches(1.5))
result_box.fill.solid()
result_box.fill.fore_color.rgb = AZUL_NAVY
result_box.line.fill.background()

txResult = slide3.shapes.add_textbox(Inches(2.5), Inches(5.0), Inches(8.333), Inches(0.5))
tfResult = txResult.text_frame
pResult = tfResult.paragraphs[0]
pResult.text = "RESULTADO"
pResult.font.size = Pt(14)
pResult.font.color.rgb = CINZA_MEDIO
pResult.alignment = PP_ALIGN.CENTER

txResult2 = slide3.shapes.add_textbox(Inches(2.5), Inches(5.4), Inches(8.333), Inches(0.6))
tfResult2 = txResult2.text_frame
pResult2 = tfResult2.paragraphs[0]
pResult2.text = "Nível de Serviço elevado + Ruptura Zero + Giro saudável"
pResult2.font.size = Pt(20)
pResult2.font.bold = True
pResult2.font.color.rgb = BRANCO
pResult2.alignment = PP_ALIGN.CENTER

# =============================================================================
# SLIDE 4 - COMO O SISTEMA FUNCIONA
# =============================================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])

bg4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg4.fill.solid()
bg4.fill.fore_color.rgb = CINZA_CLARO
bg4.line.fill.background()

add_header_bar(slide4, "COMO O SISTEMA GERA PEDIDOS INTELIGENTES", "Pipeline de cálculo automatizado")

# Fluxo de processo com 6 etapas
etapas = [
    ("1", "Histórico", "2 anos de vendas", AZUL_ACCENT),
    ("2", "Saneamento", "Correção de rupturas", VERDE_SUCESSO),
    ("3", "Previsão", "6 métodos estatísticos", LARANJA_ATENCAO),
    ("4", "Cobertura", "Lead time + ciclo + segurança", ROXO),
    ("5", "Validação", "Curva ABC e múltiplo caixa", AZUL_ACCENT),
    ("6", "Pedido", "Quantidade otimizada", VERDE_SUCESSO)
]

y_flow = 2.3
for i, (num, titulo, desc, cor) in enumerate(etapas):
    x = 0.8 + i * 2.1

    # Caixa da etapa
    box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y_flow), Inches(1.8), Inches(1.6))
    box.fill.solid()
    box.fill.fore_color.rgb = cor
    box.line.fill.background()

    # Número
    txNum = slide4.shapes.add_textbox(Inches(x), Inches(y_flow + 0.1), Inches(1.8), Inches(0.5))
    tfNum = txNum.text_frame
    pNum = tfNum.paragraphs[0]
    pNum.text = num
    pNum.font.size = Pt(24)
    pNum.font.bold = True
    pNum.font.color.rgb = BRANCO
    pNum.alignment = PP_ALIGN.CENTER

    # Título
    txT = slide4.shapes.add_textbox(Inches(x), Inches(y_flow + 0.55), Inches(1.8), Inches(0.4))
    tfT = txT.text_frame
    pT = tfT.paragraphs[0]
    pT.text = titulo
    pT.font.size = Pt(12)
    pT.font.bold = True
    pT.font.color.rgb = BRANCO
    pT.alignment = PP_ALIGN.CENTER

    # Descrição
    txD = slide4.shapes.add_textbox(Inches(x + 0.05), Inches(y_flow + 0.9), Inches(1.7), Inches(0.6))
    tfD = txD.text_frame
    tfD.word_wrap = True
    pD = tfD.paragraphs[0]
    pD.text = desc
    pD.font.size = Pt(9)
    pD.font.color.rgb = RGBColor(220, 230, 240)
    pD.alignment = PP_ALIGN.CENTER

    # Seta entre etapas
    if i < 5:
        arrow = slide4.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 1.85), Inches(y_flow + 0.65), Inches(0.2), Inches(0.25))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = CINZA_MEDIO
        arrow.line.fill.background()

# Box explicativo inferior
exp_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.3), Inches(12.333), Inches(2.5))
exp_box.fill.solid()
exp_box.fill.fore_color.rgb = BRANCO
exp_box.line.color.rgb = AZUL_ACCENT
exp_box.line.width = Pt(2)

txExp = slide4.shapes.add_textbox(Inches(0.7), Inches(4.5), Inches(12), Inches(0.4))
tfExp = txExp.text_frame
pExp = tfExp.paragraphs[0]
pExp.text = "MÉTODOS DE PREVISÃO UTILIZADOS"
pExp.font.size = Pt(14)
pExp.font.bold = True
pExp.font.color.rgb = AZUL_NAVY

# Lista de métodos em 2 colunas
metodos_col1 = ["• SMA (Média Móvel Simples) - demanda estável",
                "• WMA (Média Móvel Ponderada) - tendência recente",
                "• EMA (Suavização Exponencial) - alta variabilidade"]
metodos_col2 = ["• Regressão Linear - tendência clara",
                "• TSB (Trigg-Leach) - demanda intermitente",
                "• Decomposição Sazonal - padrão repetido"]

txM1 = slide4.shapes.add_textbox(Inches(0.7), Inches(5.0), Inches(6), Inches(1.5))
tfM1 = txM1.text_frame
for i, m in enumerate(metodos_col1):
    if i == 0:
        p = tfM1.paragraphs[0]
    else:
        p = tfM1.add_paragraph()
    p.text = m
    p.font.size = Pt(11)
    p.font.color.rgb = CINZA_TEXTO
    p.space_after = Pt(4)

txM2 = slide4.shapes.add_textbox(Inches(6.7), Inches(5.0), Inches(6), Inches(1.5))
tfM2 = txM2.text_frame
for i, m in enumerate(metodos_col2):
    if i == 0:
        p = tfM2.paragraphs[0]
    else:
        p = tfM2.add_paragraph()
    p.text = m
    p.font.size = Pt(11)
    p.font.color.rgb = CINZA_TEXTO
    p.space_after = Pt(4)

# =============================================================================
# SLIDE 5 - PARAMETRIZAÇÃO POR FORNECEDOR
# =============================================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])

bg5 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg5.fill.solid()
bg5.fill.fore_color.rgb = CINZA_CLARO
bg5.line.fill.background()

add_header_bar(slide5, "PARAMETRIZAÇÃO PERSONALIZADA", "Configurações específicas para cada fornecedor")

# Card principal de parâmetros
param_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.6), Inches(6.2), Inches(5.2))
param_box.fill.solid()
param_box.fill.fore_color.rgb = BRANCO
param_box.line.color.rgb = AZUL_ACCENT
param_box.line.width = Pt(2)

txParam = slide5.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(5.8), Inches(0.4))
tfParam = txParam.text_frame
pParam = tfParam.paragraphs[0]
pParam.text = "PARÂMETROS CONFIGURÁVEIS"
pParam.font.size = Pt(16)
pParam.font.bold = True
pParam.font.color.rgb = AZUL_NAVY

# Tabela de parâmetros
params_list = [
    ("Lead Time", "Dias entre pedido e entrega", "Configurável por loja"),
    ("Ciclo de Pedido", "Frequência de reposição", "Padrão: 7 dias"),
    ("Faturamento Mínimo", "Valor mínimo por pedido", "Define viabilidade"),
    ("Múltiplo de Caixa", "Arredondamento de quantidade", "Otimiza logística"),
    ("Nível de Serviço", "% de atendimento esperado", "98% curva A, 95% B, 90% C"),
    ("Estoque Segurança", "Dias adicionais por curva", "+2d A, +4d B, +6d C")
]

y_start = 2.3
for i, (param, desc, valor) in enumerate(params_list):
    y = y_start + i * 0.7

    # Nome do parâmetro
    txP = slide5.shapes.add_textbox(Inches(0.7), Inches(y), Inches(2.2), Inches(0.3))
    tfP = txP.text_frame
    pP = tfP.paragraphs[0]
    pP.text = param
    pP.font.size = Pt(12)
    pP.font.bold = True
    pP.font.color.rgb = AZUL_ACCENT

    # Descrição
    txD = slide5.shapes.add_textbox(Inches(0.7), Inches(y + 0.25), Inches(2.8), Inches(0.3))
    tfD = txD.text_frame
    pD = tfD.paragraphs[0]
    pD.text = desc
    pD.font.size = Pt(10)
    pD.font.color.rgb = CINZA_TEXTO

    # Valor
    txV = slide5.shapes.add_textbox(Inches(3.6), Inches(y + 0.1), Inches(2.8), Inches(0.35))
    tfV = txV.text_frame
    pV = tfV.paragraphs[0]
    pV.text = valor
    pV.font.size = Pt(10)
    pV.font.color.rgb = VERDE_SUCESSO

# Exemplo visual (direita)
ex_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.6), Inches(5.833), Inches(5.2))
ex_box.fill.solid()
ex_box.fill.fore_color.rgb = AZUL_NAVY
ex_box.line.fill.background()

txEx = slide5.shapes.add_textbox(Inches(7.2), Inches(1.8), Inches(5.4), Inches(0.4))
tfEx = txEx.text_frame
pEx = tfEx.paragraphs[0]
pEx.text = "EXEMPLO DE CÁLCULO"
pEx.font.size = Pt(14)
pEx.font.bold = True
pEx.font.color.rgb = BRANCO

# Fórmula
formula_items = [
    "Produto Curva A com Lead Time = 15 dias",
    "",
    "Cobertura Total = Lead Time + Ciclo + Segurança",
    "Cobertura Total = 15 + 7 + 2 = 24 dias",
    "",
    "Se demanda média = 100 un/dia",
    "Pedido sugerido = 24 × 100 = 2.400 unidades",
    "",
    "Arredondamento para caixa de 50:",
    "Pedido final = 2.400 unidades (48 caixas)"
]

y_ex = 2.4
for item in formula_items:
    txItem = slide5.shapes.add_textbox(Inches(7.4), Inches(y_ex), Inches(5.2), Inches(0.35))
    tfItem = txItem.text_frame
    pItem = tfItem.paragraphs[0]
    pItem.text = item
    if "=" in item and "Cobertura" not in item:
        pItem.font.size = Pt(14)
        pItem.font.bold = True
        pItem.font.color.rgb = VERDE_SUCESSO
    else:
        pItem.font.size = Pt(11)
        pItem.font.color.rgb = RGBColor(200, 210, 220)
    y_ex += 0.4

# =============================================================================
# SLIDE 6 - BALANCEAMENTO DE ESTOQUE
# =============================================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])

bg6 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg6.fill.solid()
bg6.fill.fore_color.rgb = CINZA_CLARO
bg6.line.fill.background()

add_header_bar(slide6, "BALANCEAMENTO INTELIGENTE DE ESTOQUE", "Distribuição otimizada entre lojas e CD")

# Diagrama de fluxo CD -> Lojas
# CD Central
cd_box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(1.8), Inches(2.333), Inches(1.2))
cd_box.fill.solid()
cd_box.fill.fore_color.rgb = AZUL_NAVY
cd_box.line.fill.background()

txCD = slide6.shapes.add_textbox(Inches(5.5), Inches(2.15), Inches(2.333), Inches(0.6))
tfCD = txCD.text_frame
pCD = tfCD.paragraphs[0]
pCD.text = "CD CENTRAL"
pCD.font.size = Pt(16)
pCD.font.bold = True
pCD.font.color.rgb = BRANCO
pCD.alignment = PP_ALIGN.CENTER

# Lojas
lojas_info = [
    ("LOJA 1", "Cobertura: 28d", VERDE_SUCESSO, 1.0),
    ("LOJA 2", "Cobertura: 5d", VERMELHO, 4.5),
    ("LOJA 3", "Cobertura: 15d", LARANJA_ATENCAO, 8.0),
    ("LOJA 4", "Cobertura: 22d", VERDE_SUCESSO, 11.5)
]

for nome, cob, cor, x in lojas_info:
    loja_box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(4.0), Inches(2.0), Inches(1.0))
    loja_box.fill.solid()
    loja_box.fill.fore_color.rgb = BRANCO
    loja_box.line.color.rgb = cor
    loja_box.line.width = Pt(3)

    txL = slide6.shapes.add_textbox(Inches(x), Inches(4.1), Inches(2.0), Inches(0.4))
    tfL = txL.text_frame
    pL = tfL.paragraphs[0]
    pL.text = nome
    pL.font.size = Pt(12)
    pL.font.bold = True
    pL.font.color.rgb = AZUL_NAVY
    pL.alignment = PP_ALIGN.CENTER

    txC = slide6.shapes.add_textbox(Inches(x), Inches(4.5), Inches(2.0), Inches(0.4))
    tfC = txC.text_frame
    pC = tfC.paragraphs[0]
    pC.text = cob
    pC.font.size = Pt(10)
    pC.font.color.rgb = cor
    pC.alignment = PP_ALIGN.CENTER

    # Linhas conectoras
    line = slide6.shapes.add_connector(1, Inches(6.666), Inches(3.0), Inches(x + 1.0), Inches(4.0))
    line.line.color.rgb = CINZA_MEDIO

# Transferência sugerida
transf_box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(5.3), Inches(8.333), Inches(1.5))
transf_box.fill.solid()
transf_box.fill.fore_color.rgb = BRANCO
transf_box.line.color.rgb = VERDE_SUCESSO
transf_box.line.width = Pt(2)

txTransf = slide6.shapes.add_textbox(Inches(2.7), Inches(5.45), Inches(8), Inches(0.4))
tfTransf = txTransf.text_frame
pTransf = tfTransf.paragraphs[0]
pTransf.text = "TRANSFERÊNCIA SUGERIDA"
pTransf.font.size = Pt(12)
pTransf.font.bold = True
pTransf.font.color.rgb = VERDE_SUCESSO

txTransf2 = slide6.shapes.add_textbox(Inches(2.7), Inches(5.85), Inches(8), Inches(0.8))
tfTransf2 = txTransf2.text_frame
pTransf2 = tfTransf2.paragraphs[0]
pTransf2.text = "Sistema identifica: LOJA 1 (excesso) → LOJA 2 (falta) = 500 unidades"
pTransf2.font.size = Pt(11)
pTransf2.font.color.rgb = CINZA_TEXTO

pTransf3 = tfTransf2.add_paragraph()
pTransf3.text = "Resultado: Equaliza cobertura, evita ruptura e reduz pedido emergencial"
pTransf3.font.size = Pt(11)
pTransf3.font.color.rgb = AZUL_ACCENT

# =============================================================================
# SLIDE 7 - PROCESSO DE GESTÃO DE ABASTECIMENTO
# =============================================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])

bg7 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg7.fill.solid()
bg7.fill.fore_color.rgb = CINZA_CLARO
bg7.line.fill.background()

add_header_bar(slide7, "PROCESSO DE GESTÃO DE ABASTECIMENTO", "Ciclo completo de pedido com rastreabilidade")

# Processo circular com 4 etapas
processo_items = [
    ("1", "PEDIDO\nCICLO FIXO", "Emissão programada\ncom frequência definida", AZUL_ACCENT, 1.5, 2.0),
    ("2", "TRACKING\nDO PEDIDO", "Acompanhamento em\ntempo real do status", VERDE_SUCESSO, 7.5, 2.0),
    ("3", "ENTREGA\nINTEGRAL", "Compromisso com\n100% do pedido", LARANJA_ATENCAO, 7.5, 4.5),
    ("4", "AVALIAÇÃO\nNÍVEL SERVIÇO", "Métricas de performance\ndo fornecedor", ROXO, 1.5, 4.5)
]

for num, titulo, desc, cor, x, y in processo_items:
    # Box principal
    box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(4.0), Inches(1.8))
    box.fill.solid()
    box.fill.fore_color.rgb = cor
    box.line.fill.background()

    # Número
    txN = slide7.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.2), Inches(0.5), Inches(0.5))
    tfN = txN.text_frame
    pN = tfN.paragraphs[0]
    pN.text = num
    pN.font.size = Pt(24)
    pN.font.bold = True
    pN.font.color.rgb = BRANCO

    # Título
    txT = slide7.shapes.add_textbox(Inches(x + 0.8), Inches(y + 0.15), Inches(3.0), Inches(0.7))
    tfT = txT.text_frame
    tfT.word_wrap = True
    pT = tfT.paragraphs[0]
    pT.text = titulo
    pT.font.size = Pt(14)
    pT.font.bold = True
    pT.font.color.rgb = BRANCO

    # Descrição
    txD = slide7.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.95), Inches(3.6), Inches(0.8))
    tfD = txD.text_frame
    tfD.word_wrap = True
    pD = tfD.paragraphs[0]
    pD.text = desc
    pD.font.size = Pt(10)
    pD.font.color.rgb = RGBColor(220, 230, 240)

# Setas de fluxo
arrow_right = slide7.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.6), Inches(2.7), Inches(1.8), Inches(0.3))
arrow_right.fill.solid()
arrow_right.fill.fore_color.rgb = CINZA_MEDIO
arrow_right.line.fill.background()

arrow_down = slide7.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(9.3), Inches(3.9), Inches(0.3), Inches(0.5))
arrow_down.fill.solid()
arrow_down.fill.fore_color.rgb = CINZA_MEDIO
arrow_down.line.fill.background()

arrow_left = slide7.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.6), Inches(5.2), Inches(1.8), Inches(0.3))
arrow_left.fill.solid()
arrow_left.fill.fore_color.rgb = CINZA_MEDIO
arrow_left.line.fill.background()
arrow_left.rotation = 180

arrow_up = slide7.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(3.3), Inches(3.9), Inches(0.3), Inches(0.5))
arrow_up.fill.solid()
arrow_up.fill.fore_color.rgb = CINZA_MEDIO
arrow_up.line.fill.background()
arrow_up.rotation = 180

# Centro - Relacionamento
center_box = slide7.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.666), Inches(3.5), Inches(2.0), Inches(1.5))
center_box.fill.solid()
center_box.fill.fore_color.rgb = AZUL_NAVY
center_box.line.fill.background()

txCenter = slide7.shapes.add_textbox(Inches(5.666), Inches(3.85), Inches(2.0), Inches(0.8))
tfCenter = txCenter.text_frame
tfCenter.word_wrap = True
pCenter = tfCenter.paragraphs[0]
pCenter.text = "PARCERIA"
pCenter.font.size = Pt(12)
pCenter.font.bold = True
pCenter.font.color.rgb = BRANCO
pCenter.alignment = PP_ALIGN.CENTER

pCenter2 = tfCenter.add_paragraph()
pCenter2.text = "PRÓXIMA"
pCenter2.font.size = Pt(12)
pCenter2.font.bold = True
pCenter2.font.color.rgb = BRANCO
pCenter2.alignment = PP_ALIGN.CENTER

# =============================================================================
# SLIDE 8 - O QUE ESPERAMOS DO FORNECEDOR
# =============================================================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])

bg8 = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg8.fill.solid()
bg8.fill.fore_color.rgb = CINZA_CLARO
bg8.line.fill.background()

add_header_bar(slide8, "O QUE ESPERAMOS DO FORNECEDOR", "Compromissos para elevar o nível de serviço")

# Cards de compromissos
compromissos = [
    ("ENTREGA INTEGRAL",
     ["100% do pedido entregue", "Sem cortes ou substituições", "Comunicação prévia de faltas"],
     VERDE_SUCESSO, "100%"),
    ("PRAZO ACORDADO",
     ["Respeito ao lead time", "Agendamento de entrega", "Pontualidade na chegada"],
     AZUL_ACCENT, "LT"),
    ("QUALIDADE",
     ["Produtos conforme pedido", "Embalagem adequada", "Documentação correta"],
     LARANJA_ATENCAO, "QA"),
    ("COMUNICAÇÃO",
     ["Aviso antecipado de problemas", "Resposta ágil a consultas", "Transparência no status"],
     ROXO, "COM")
]

x_positions = [0.5, 3.5, 6.5, 9.5]
for i, (titulo, items, cor, icone) in enumerate(compromissos):
    add_card(slide8, x_positions[i], 1.6, 2.8, 3.5, titulo, items, cor, icone)

# Box de impacto
impact_box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.4), Inches(12.333), Inches(1.4))
impact_box.fill.solid()
impact_box.fill.fore_color.rgb = AZUL_NAVY
impact_box.line.fill.background()

txImpact = slide8.shapes.add_textbox(Inches(0.7), Inches(5.55), Inches(12), Inches(0.4))
tfImpact = txImpact.text_frame
pImpact = tfImpact.paragraphs[0]
pImpact.text = "IMPACTO POSITIVO PARA AMBOS"
pImpact.font.size = Pt(14)
pImpact.font.bold = True
pImpact.font.color.rgb = VERDE_SUCESSO

txImpact2 = slide8.shapes.add_textbox(Inches(0.7), Inches(5.95), Inches(12), Inches(0.7))
tfImpact2 = txImpact2.text_frame
tfImpact2.word_wrap = True
pImpact2 = tfImpact2.paragraphs[0]
pImpact2.text = "Fornecedores com alto nível de serviço terão prioridade nas compras, volumes maiores e relacionamento de longo prazo."
pImpact2.font.size = Pt(13)
pImpact2.font.color.rgb = BRANCO

# =============================================================================
# SLIDE 9 - MÉTRICAS E ACOMPANHAMENTO
# =============================================================================
slide9 = prs.slides.add_slide(prs.slide_layouts[6])

bg9 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg9.fill.solid()
bg9.fill.fore_color.rgb = CINZA_CLARO
bg9.line.fill.background()

add_header_bar(slide9, "MÉTRICAS DE ACOMPANHAMENTO", "Como medimos o nível de serviço do fornecedor")

# KPIs principais
add_metric_box(slide9, 0.5, 1.8, "98%", "Nível de Serviço\nCurva A", VERDE_SUCESSO, 2.5)
add_metric_box(slide9, 3.3, 1.8, "95%", "Nível de Serviço\nCurva B", AZUL_ACCENT, 2.5)
add_metric_box(slide9, 6.1, 1.8, "90%", "Nível de Serviço\nCurva C", LARANJA_ATENCAO, 2.5)
add_metric_box(slide9, 8.9, 1.8, "<5%", "Taxa de\nRuptura", ROXO, 2.5)

# Tabela de métricas detalhadas
metricas_box = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.5), Inches(12.333), Inches(3.3))
metricas_box.fill.solid()
metricas_box.fill.fore_color.rgb = BRANCO
metricas_box.line.color.rgb = AZUL_ACCENT
metricas_box.line.width = Pt(2)

txMet = slide9.shapes.add_textbox(Inches(0.7), Inches(3.7), Inches(12), Inches(0.4))
tfMet = txMet.text_frame
pMet = tfMet.paragraphs[0]
pMet.text = "INDICADORES MONITORADOS"
pMet.font.size = Pt(14)
pMet.font.bold = True
pMet.font.color.rgb = AZUL_NAVY

metricas = [
    ("Fill Rate", "% do pedido entregue completo", "Meta: ≥ 95%", VERDE_SUCESSO),
    ("On-Time Delivery", "% de entregas no prazo acordado", "Meta: ≥ 90%", AZUL_ACCENT),
    ("Lead Time Real", "Dias efetivos vs. acordado", "Meta: ≤ LT cadastrado", LARANJA_ATENCAO),
    ("Qualidade", "% de entregas sem problemas", "Meta: ≥ 98%", ROXO)
]

y_met = 4.2
for metrica, desc, meta, cor in metricas:
    # Indicador colorido
    ind = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(y_met), Inches(0.15), Inches(0.5))
    ind.fill.solid()
    ind.fill.fore_color.rgb = cor
    ind.line.fill.background()

    # Nome
    txN = slide9.shapes.add_textbox(Inches(1.0), Inches(y_met), Inches(2.5), Inches(0.3))
    tfN = txN.text_frame
    pN = tfN.paragraphs[0]
    pN.text = metrica
    pN.font.size = Pt(12)
    pN.font.bold = True
    pN.font.color.rgb = AZUL_NAVY

    # Descrição
    txD = slide9.shapes.add_textbox(Inches(3.5), Inches(y_met + 0.05), Inches(5.5), Inches(0.4))
    tfD = txD.text_frame
    pD = tfD.paragraphs[0]
    pD.text = desc
    pD.font.size = Pt(11)
    pD.font.color.rgb = CINZA_TEXTO

    # Meta
    txM = slide9.shapes.add_textbox(Inches(9.5), Inches(y_met + 0.05), Inches(3), Inches(0.4))
    tfM = txM.text_frame
    pM = tfM.paragraphs[0]
    pM.text = meta
    pM.font.size = Pt(11)
    pM.font.bold = True
    pM.font.color.rgb = cor

    y_met += 0.6

# =============================================================================
# SLIDE 10 - BENEFÍCIOS DA PARCERIA
# =============================================================================
slide10 = prs.slides.add_slide(prs.slide_layouts[6])

bg10 = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg10.fill.solid()
bg10.fill.fore_color.rgb = CINZA_CLARO
bg10.line.fill.background()

add_header_bar(slide10, "BENEFÍCIOS DA PARCERIA", "Ganhos mútuos com a gestão inteligente")

# Lado fornecedor
forn_box = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.6), Inches(5.9), Inches(5.2))
forn_box.fill.solid()
forn_box.fill.fore_color.rgb = VERDE_SUCESSO
forn_box.line.fill.background()

txForn = slide10.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(5.5), Inches(0.5))
tfForn = txForn.text_frame
pForn = tfForn.paragraphs[0]
pForn.text = "PARA O FORNECEDOR"
pForn.font.size = Pt(18)
pForn.font.bold = True
pForn.font.color.rgb = BRANCO

beneficios_forn = [
    "Previsibilidade de demanda para planejamento",
    "Pedidos regulares com ciclo fixo",
    "Redução de devoluções e trocas",
    "Relacionamento de longo prazo",
    "Prioridade para fornecedores bem avaliados",
    "Comunicação transparente e objetiva",
    "Oportunidade de crescimento conjunto"
]

y_bf = 2.5
for ben in beneficios_forn:
    txB = slide10.shapes.add_textbox(Inches(0.9), Inches(y_bf), Inches(5.3), Inches(0.45))
    tfB = txB.text_frame
    tfB.word_wrap = True
    pB = tfB.paragraphs[0]
    pB.text = f"✓ {ben}"
    pB.font.size = Pt(12)
    pB.font.color.rgb = BRANCO
    y_bf += 0.55

# Lado Ferreira Costa
fc_box = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.933), Inches(1.6), Inches(5.9), Inches(5.2))
fc_box.fill.solid()
fc_box.fill.fore_color.rgb = AZUL_ACCENT
fc_box.line.fill.background()

txFC = slide10.shapes.add_textbox(Inches(7.133), Inches(1.8), Inches(5.5), Inches(0.5))
tfFC = txFC.text_frame
pFC = tfFC.paragraphs[0]
pFC.text = "PARA FERREIRA COSTA"
pFC.font.size = Pt(18)
pFC.font.bold = True
pFC.font.color.rgb = BRANCO

beneficios_fc = [
    "Redução de rupturas nas gôndolas",
    "Estoque balanceado e otimizado",
    "Menor capital imobilizado",
    "Satisfação do cliente final",
    "Eficiência operacional",
    "Dados para decisões estratégicas",
    "Competitividade no mercado"
]

y_bfc = 2.5
for ben in beneficios_fc:
    txB = slide10.shapes.add_textbox(Inches(7.333), Inches(y_bfc), Inches(5.3), Inches(0.45))
    tfB = txB.text_frame
    tfB.word_wrap = True
    pB = tfB.paragraphs[0]
    pB.text = f"✓ {ben}"
    pB.font.size = Pt(12)
    pB.font.color.rgb = BRANCO
    y_bfc += 0.55

# =============================================================================
# SLIDE 11 - PRÓXIMOS PASSOS
# =============================================================================
slide11 = prs.slides.add_slide(prs.slide_layouts[6])

bg11 = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg11.fill.solid()
bg11.fill.fore_color.rgb = CINZA_CLARO
bg11.line.fill.background()

add_header_bar(slide11, "PRÓXIMOS PASSOS", "Implantação do processo de gestão de abastecimento")

# Timeline
passos = [
    ("1", "Cadastro de\nParâmetros", "Definição de lead time,\nciclo e faturamento mínimo", AZUL_ACCENT),
    ("2", "Validação\nConjunta", "Revisão dos parâmetros\ne ajustes necessários", VERDE_SUCESSO),
    ("3", "Início do\nCiclo Fixo", "Primeiro pedido com\nnovas regras", LARANJA_ATENCAO),
    ("4", "Acompanhamento\nMensal", "Revisão de métricas\ne ajustes finos", ROXO)
]

x_start = 0.8
for i, (num, titulo, desc, cor) in enumerate(passos):
    x = x_start + i * 3.15

    # Círculo numerado
    circle = slide11.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.85), Inches(2.0), Inches(0.9), Inches(0.9))
    circle.fill.solid()
    circle.fill.fore_color.rgb = cor
    circle.line.fill.background()

    txNum = slide11.shapes.add_textbox(Inches(x + 0.85), Inches(2.2), Inches(0.9), Inches(0.6))
    tfNum = txNum.text_frame
    pNum = tfNum.paragraphs[0]
    pNum.text = num
    pNum.font.size = Pt(28)
    pNum.font.bold = True
    pNum.font.color.rgb = BRANCO
    pNum.alignment = PP_ALIGN.CENTER

    # Box de conteúdo
    box = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.1), Inches(2.7), Inches(2.0))
    box.fill.solid()
    box.fill.fore_color.rgb = BRANCO
    box.line.color.rgb = cor
    box.line.width = Pt(2)

    # Título
    txT = slide11.shapes.add_textbox(Inches(x + 0.1), Inches(3.2), Inches(2.5), Inches(0.7))
    tfT = txT.text_frame
    tfT.word_wrap = True
    pT = tfT.paragraphs[0]
    pT.text = titulo
    pT.font.size = Pt(13)
    pT.font.bold = True
    pT.font.color.rgb = AZUL_NAVY
    pT.alignment = PP_ALIGN.CENTER

    # Descrição
    txD = slide11.shapes.add_textbox(Inches(x + 0.1), Inches(3.9), Inches(2.5), Inches(1.0))
    tfD = txD.text_frame
    tfD.word_wrap = True
    pD = tfD.paragraphs[0]
    pD.text = desc
    pD.font.size = Pt(10)
    pD.font.color.rgb = CINZA_TEXTO
    pD.alignment = PP_ALIGN.CENTER

    # Linha conectora
    if i < 3:
        line = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 2.8), Inches(2.4), Inches(0.35), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = CINZA_MEDIO
        line.line.fill.background()

# Box de ação
action_box = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(5.5), Inches(8.333), Inches(1.3))
action_box.fill.solid()
action_box.fill.fore_color.rgb = AZUL_NAVY
action_box.line.fill.background()

txAction = slide11.shapes.add_textbox(Inches(2.7), Inches(5.7), Inches(8), Inches(0.4))
tfAction = txAction.text_frame
pAction = tfAction.paragraphs[0]
pAction.text = "AÇÃO IMEDIATA"
pAction.font.size = Pt(12)
pAction.font.bold = True
pAction.font.color.rgb = VERDE_SUCESSO
pAction.alignment = PP_ALIGN.CENTER

txAction2 = slide11.shapes.add_textbox(Inches(2.7), Inches(6.1), Inches(8), Inches(0.5))
tfAction2 = txAction2.text_frame
pAction2 = tfAction2.paragraphs[0]
pAction2.text = "Preencher planilha de parâmetros e enviar para validação"
pAction2.font.size = Pt(14)
pAction2.font.color.rgb = BRANCO
pAction2.alignment = PP_ALIGN.CENTER

# =============================================================================
# SLIDE 12 - ENCERRAMENTO
# =============================================================================
slide12 = prs.slides.add_slide(prs.slide_layouts[6])

# Fundo
bg12 = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg12.fill.solid()
bg12.fill.fore_color.rgb = AZUL_NAVY
bg12.line.fill.background()

# Linha decorativa
line12 = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3.5), Inches(12.333), Inches(0.02))
line12.fill.solid()
line12.fill.fore_color.rgb = AZUL_ACCENT
line12.line.fill.background()

# Título
txTitle12 = slide12.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(12.333), Inches(1.0))
tf12 = txTitle12.text_frame
p12 = tf12.paragraphs[0]
p12.text = "JUNTOS SOMOS MAIS FORTES"
p12.font.size = Pt(44)
p12.font.bold = True
p12.font.color.rgb = BRANCO
p12.alignment = PP_ALIGN.CENTER

# Subtítulo
txSub12 = slide12.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(12.333), Inches(0.6))
tf12b = txSub12.text_frame
p12b = tf12b.paragraphs[0]
p12b.text = "Conte conosco para construir uma cadeia de suprimentos eficiente"
p12b.font.size = Pt(20)
p12b.font.color.rgb = CINZA_MEDIO
p12b.alignment = PP_ALIGN.CENTER

# Métricas de impacto
metricas_final = [
    ("Redução de Rupturas", VERDE_SUCESSO),
    ("Aumento de Vendas", AZUL_ACCENT),
    ("Parceria Duradoura", LARANJA_ATENCAO)
]

x_met = 2.5
for metrica, cor in metricas_final:
    box = slide12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_met), Inches(4.5), Inches(2.5), Inches(0.8))
    box.fill.solid()
    box.fill.fore_color.rgb = cor
    box.line.fill.background()

    txM = slide12.shapes.add_textbox(Inches(x_met), Inches(4.7), Inches(2.5), Inches(0.5))
    tfM = txM.text_frame
    pM = tfM.paragraphs[0]
    pM.text = metrica
    pM.font.size = Pt(12)
    pM.font.bold = True
    pM.font.color.rgb = BRANCO
    pM.alignment = PP_ALIGN.CENTER

    x_met += 2.9

# Contato
txContact = slide12.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12.333), Inches(0.5))
tfContact = txContact.text_frame
pContact = tfContact.paragraphs[0]
pContact.text = "Equipe de Abastecimento - Ferreira Costa"
pContact.font.size = Pt(14)
pContact.font.color.rgb = CINZA_MEDIO
pContact.alignment = PP_ALIGN.CENTER

# Logo
txLogo12 = slide12.shapes.add_textbox(Inches(5.5), Inches(6.5), Inches(2.333), Inches(0.5))
tfLogo12 = txLogo12.text_frame
pLogo12 = tfLogo12.paragraphs[0]
pLogo12.text = "FERREIRA COSTA"
pLogo12.font.size = Pt(16)
pLogo12.font.bold = True
pLogo12.font.color.rgb = BRANCO
pLogo12.alignment = PP_ALIGN.CENTER

# =============================================================================
# SALVAR APRESENTAÇÃO
# =============================================================================
output_path = r"C:\temp\Apresentacao_Fornecedores_v1.pptx"
prs.save(output_path)

print(f"\nApresentação criada com sucesso!")
print(f"Arquivo: {output_path}")
print(f"\nTotal de slides: {len(prs.slides)}")
print("\nConteúdo dos slides:")
print("  1. Capa - Parceria para Excelência")
print("  2. Nosso Objetivo Comum - Diagrama de benefícios")
print("  3. O que é um Pedido de Qualidade - 4 Pilares")
print("  4. Como o Sistema Funciona - Pipeline de cálculo")
print("  5. Parametrização Personalizada - Configurações por fornecedor")
print("  6. Balanceamento de Estoque - CD e Lojas")
print("  7. Processo de Gestão de Abastecimento - Ciclo completo")
print("  8. O que Esperamos do Fornecedor - Compromissos")
print("  9. Métricas de Acompanhamento - KPIs")
print(" 10. Benefícios da Parceria - Ganhos mútuos")
print(" 11. Próximos Passos - Timeline de implantação")
print(" 12. Encerramento - Juntos somos mais fortes")
