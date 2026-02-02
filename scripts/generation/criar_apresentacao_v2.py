# -*- coding: utf-8 -*-
"""Script para criar apresentação PowerPoint - Versão Visual/Executiva"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import nsmap
import os

# Criar apresentação
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Cores
AZUL_ESCURO = RGBColor(0, 51, 102)
AZUL_CLARO = RGBColor(0, 112, 192)
AZUL_LIGHT = RGBColor(217, 226, 243)
VERDE = RGBColor(0, 176, 80)
VERMELHO = RGBColor(192, 0, 0)
LARANJA = RGBColor(255, 153, 0)
CINZA = RGBColor(127, 127, 127)
BRANCO = RGBColor(255, 255, 255)
PRETO = RGBColor(51, 51, 51)

def add_title_slide(title, subtitle=""):
    """Slide de título com design moderno"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Fundo gradiente simulado com retângulos
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = AZUL_ESCURO
    shape.line.fill.background()

    # Elemento decorativo - linha lateral
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(2), Inches(0.08), Inches(3))
    line.fill.solid()
    line.fill.fore_color.rgb = LARANJA
    line.line.fill.background()

    # Título
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(11), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = BRANCO

    if subtitle:
        p = tf.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.space_before = Pt(20)

    return slide

def add_header(slide, title):
    """Adiciona header padrão ao slide"""
    # Barra superior
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = AZUL_ESCURO
    shape.line.fill.background()

    # Título
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BRANCO

def add_icon_box(slide, x, y, width, height, icon_text, label, color, desc=""):
    """Adiciona caixa com ícone estilizado"""
    # Círculo/quadrado com ícone
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

    # Texto do ícone (emoji ou símbolo)
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y + 0.15), Inches(width), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = icon_text
    p.font.size = Pt(32)
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER

    # Label abaixo
    txBox = slide.shapes.add_textbox(Inches(x - 0.3), Inches(y + height + 0.1), Inches(width + 0.6), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRETO
    p.alignment = PP_ALIGN.CENTER

    if desc:
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = CINZA
        p.alignment = PP_ALIGN.CENTER

def add_arrow(slide, x1, y1, x2, y2):
    """Adiciona seta entre elementos"""
    connector = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x1), Inches(y1), Inches(x2-x1), Inches(0.3))
    connector.fill.solid()
    connector.fill.fore_color.rgb = AZUL_CLARO
    connector.line.fill.background()

def add_kpi_card(slide, x, y, value, label, color):
    """Adiciona card de KPI"""
    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.8), Inches(1.5))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(250, 250, 250)
    card.line.color.rgb = color
    card.line.width = Pt(3)

    # Valor grande
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y + 0.2), Inches(2.8), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER

    # Label
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y + 0.95), Inches(2.8), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(12)
    p.font.color.rgb = CINZA
    p.alignment = PP_ALIGN.CENTER

def add_comparison_box(slide, x, y, title, items, is_before=True):
    """Adiciona caixa de comparação Antes/Depois"""
    color = VERMELHO if is_before else VERDE
    bg_color = RGBColor(255, 235, 235) if is_before else RGBColor(235, 255, 235)

    # Background
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.5), Inches(4.5))
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    box.line.color.rgb = color
    box.line.width = Pt(2)

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.5), Inches(0.7))
    header.fill.solid()
    header.fill.fore_color.rgb = color
    header.line.fill.background()

    # Título
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y + 0.15), Inches(5.5), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER

    # Items com ícones
    y_item = y + 0.9
    icon = "X" if is_before else "+"
    for item in items:
        txBox = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y_item), Inches(5), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"{icon}  {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = PRETO
        y_item += 0.55

def add_timeline_item(slide, x, y, number, title, active=False):
    """Adiciona item de timeline/roadmap"""
    color = AZUL_CLARO if active else CINZA

    # Círculo numerado
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()

    # Número
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y + 0.15), Inches(0.7), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER

    # Título abaixo
    txBox = slide.shapes.add_textbox(Inches(x - 0.5), Inches(y + 0.85), Inches(1.7), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER

# ============ CRIAR SLIDES ============

# SLIDE 1 - Capa
add_title_slide(
    "PREVISÃO DE DEMANDA",
    "Proposta de Implantação Piloto  |  Janeiro 2026"
)

# SLIDE 2 - O Desafio (visual com ícones)
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "O DESAFIO ATUAL")

# Ícones representando problemas
add_icon_box(slide, 1, 2, 1.2, 1.2, "!", "RUPTURA", VERMELHO, "Perda de vendas")
add_icon_box(slide, 3.5, 2, 1.2, 1.2, "$", "ESTOQUE", VERMELHO, "Capital parado")
add_icon_box(slide, 6, 2, 1.2, 1.2, "?", "DECISÃO", VERMELHO, "Sem dados")
add_icon_box(slide, 8.5, 2, 1.2, 1.2, "T", "TEMPO", VERMELHO, "Processo manual")
add_icon_box(slide, 11, 2, 1.2, 1.2, "X", "ERRO", VERMELHO, "Pedidos imprecisos")

# Texto central
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Hoje: Cálculo manual de pedidos pela equipe comercial"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = PRETO
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "Mesma regra para todos os itens e lojas, sem considerar sazonalidade ou histórico"
p.font.size = Pt(18)
p.font.color.rgb = CINZA
p.alignment = PP_ALIGN.CENTER

# SLIDE 3 - A Solução (fluxo visual)
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "A SOLUÇÃO")

# Fluxo: Dados -> Sistema -> Resultados
# Box Dados
box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.5), Inches(3), Inches(3))
box1.fill.solid()
box1.fill.fore_color.rgb = AZUL_LIGHT
box1.line.color.rgb = AZUL_ESCURO
box1.line.width = Pt(2)

txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.7), Inches(3), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "DADOS"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = AZUL_ESCURO
p.alignment = PP_ALIGN.CENTER

txBox = slide.shapes.add_textbox(Inches(0.7), Inches(3.3), Inches(2.6), Inches(2))
tf = txBox.text_frame
for item in ["3 anos de histórico", "10M+ registros", "Estoque atual", "Cadastro produtos"]:
    p = tf.add_paragraph()
    p.text = f"• {item}"
    p.font.size = Pt(14)
    p.font.color.rgb = PRETO

# Seta 1
add_arrow(slide, 3.7, 3.8, 4.8, 3.8)

# Box Sistema (central, maior)
box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5), Inches(2), Inches(3.5), Inches(4))
box2.fill.solid()
box2.fill.fore_color.rgb = AZUL_ESCURO
box2.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(5), Inches(2.2), Inches(3.5), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "SISTEMA"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

txBox = slide.shapes.add_textbox(Inches(5.2), Inches(2.8), Inches(3.1), Inches(3))
tf = txBox.text_frame
for item in ["Algoritmos estatísticos", "Machine Learning", "Sazonalidade", "Tendências"]:
    p = tf.add_paragraph()
    p.text = f"• {item}"
    p.font.size = Pt(14)
    p.font.color.rgb = BRANCO
    p.space_after = Pt(8)

# Seta 2
add_arrow(slide, 8.7, 3.8, 9.8, 3.8)

# Box Resultados
box3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10), Inches(2.5), Inches(3), Inches(3))
box3.fill.solid()
box3.fill.fore_color.rgb = RGBColor(235, 255, 235)
box3.line.color.rgb = VERDE
box3.line.width = Pt(2)

txBox = slide.shapes.add_textbox(Inches(10), Inches(2.7), Inches(3), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "RESULTADOS"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = VERDE
p.alignment = PP_ALIGN.CENTER

txBox = slide.shapes.add_textbox(Inches(10.2), Inches(3.3), Inches(2.6), Inches(2))
tf = txBox.text_frame
for item in ["Previsão 12 meses", "Pedidos otimizados", "Menos ruptura", "Mais vendas"]:
    p = tf.add_paragraph()
    p.text = f"• {item}"
    p.font.size = Pt(14)
    p.font.color.rgb = PRETO

# SLIDE 4 - Módulos do Sistema (4 quadrantes)
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "MÓDULOS DO SISTEMA")

modules = [
    (0.5, 1.5, "PREVISÃO", AZUL_CLARO, ["Forecast 12 meses", "Por loja/produto", "Sazonalidade"]),
    (6.9, 1.5, "PEDIDOS", VERDE, ["Geração automática", "Arredondamento", "Calendário"]),
    (0.5, 4.2, "ESTOQUE", LARANJA, ["Posição atual", "Alertas ruptura", "Balanceamento"]),
    (6.9, 4.2, "FORNECEDORES", RGBColor(128, 0, 128), ["Nível serviço", "Tracking", "Performance"])
]

for x, y, title, color, items in modules:
    # Box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(6), Inches(2.5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(250, 250, 250)
    box.line.color.rgb = color
    box.line.width = Pt(3)

    # Barra colorida lateral
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.15), Inches(2.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    # Título
    txBox = slide.shapes.add_textbox(Inches(x + 0.4), Inches(y + 0.2), Inches(5.5), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = color

    # Items
    txBox = slide.shapes.add_textbox(Inches(x + 0.4), Inches(y + 0.8), Inches(5.5), Inches(1.5))
    tf = txBox.text_frame
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = PRETO

# SLIDE 5 - Antes vs Depois
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "TRANSFORMAÇÃO")

add_comparison_box(slide, 0.5, 1.5, "HOJE", [
    "Cálculo manual",
    "Sem previsão estruturada",
    "Decisão por intuição",
    "Mesma regra para tudo",
    "Sem visibilidade",
    "Reativo"
], is_before=True)

# Seta central grande
arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.2), Inches(3.5), Inches(1), Inches(0.8))
arrow.fill.solid()
arrow.fill.fore_color.rgb = AZUL_CLARO
arrow.line.fill.background()

add_comparison_box(slide, 7.3, 1.5, "COM O SISTEMA", [
    "Cálculo automático",
    "Previsão 12 meses",
    "Decisão baseada em dados",
    "Personalizado por item/loja",
    "Dashboards em tempo real",
    "Proativo"
], is_before=False)

# SLIDE 6 - Proposta Piloto (visual com números)
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "PROPOSTA: PILOTO CONTROLADO")

# Cards com números grandes
add_kpi_card(slide, 0.5, 1.8, "4", "FORNECEDORES", AZUL_CLARO)
add_kpi_card(slide, 3.5, 1.8, "701", "PRODUTOS", VERDE)
add_kpi_card(slide, 6.5, 1.8, "10M+", "REGISTROS", LARANJA)
add_kpi_card(slide, 9.5, 1.8, "3", "ANOS HISTÓRICO", RGBColor(128, 0, 128))

# Lista de fornecedores visual
fornecedores = [
    ("FAME", "45 SKUs", "Ferragens"),
    ("ZAGONEL", "75 SKUs", "Chuveiros"),
    ("SILVANA", "311 SKUs", "Fechaduras"),
    ("STAM", "270 SKUs", "Fechaduras")
]

y = 4
for i, (nome, skus, cat) in enumerate(fornecedores):
    x = 0.5 + i * 3.2
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3), Inches(1.8))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(250, 250, 250)
    box.line.color.rgb = AZUL_CLARO
    box.line.width = Pt(2)

    txBox = slide.shapes.add_textbox(Inches(x), Inches(y + 0.2), Inches(3), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = nome
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = AZUL_ESCURO
    p.alignment = PP_ALIGN.CENTER

    txBox = slide.shapes.add_textbox(Inches(x), Inches(y + 0.7), Inches(3), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = skus
    p.font.size = Pt(14)
    p.font.color.rgb = CINZA
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = cat
    p.font.size = Pt(12)
    p.font.color.rgb = CINZA
    p.alignment = PP_ALIGN.CENTER

# Texto
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(12), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Todas as lojas da rede  |  Dados já carregados e validados  |  Sistema pronto"
p.font.size = Pt(16)
p.font.color.rgb = CINZA
p.alignment = PP_ALIGN.CENTER

# SLIDE 7 - Benefícios (KPIs visuais)
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "BENEFÍCIOS ESPERADOS")

beneficios = [
    ("-30%", "RUPTURA", VERDE, "Mais vendas"),
    ("+40%", "ACURACIDADE", AZUL_CLARO, "Pedidos precisos"),
    ("+25%", "EFICIÊNCIA", LARANJA, "Estoque otimizado"),
    ("+50%", "PRODUTIVIDADE", RGBColor(128, 0, 128), "Menos trabalho manual")
]

for i, (valor, titulo, cor, desc) in enumerate(beneficios):
    x = 0.5 + i * 3.2

    # Card grande
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.8), Inches(3), Inches(3.5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(250, 250, 250)
    box.line.color.rgb = cor
    box.line.width = Pt(4)

    # Valor grande
    txBox = slide.shapes.add_textbox(Inches(x), Inches(2.2), Inches(3), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = valor
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = cor
    p.alignment = PP_ALIGN.CENTER

    # Título
    txBox = slide.shapes.add_textbox(Inches(x), Inches(3.3), Inches(3), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = PRETO
    p.alignment = PP_ALIGN.CENTER

    # Descrição
    txBox = slide.shapes.add_textbox(Inches(x), Inches(3.9), Inches(3), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(14)
    p.font.color.rgb = CINZA
    p.alignment = PP_ALIGN.CENTER

# Impactos adicionais
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12), Inches(1))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Impacto Financeiro: Redução de custos logísticos | Otimização do capital de giro | Melhor fluxo de caixa"
p.font.size = Pt(16)
p.font.color.rgb = AZUL_ESCURO
p.alignment = PP_ALIGN.CENTER

# SLIDE 8 - Roadmap (timeline visual)
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "ROADMAP DE EXPANSÃO")

# Linha de conexão
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(3.05), Inches(10.5), Inches(0.1))
line.fill.solid()
line.fill.fore_color.rgb = RGBColor(200, 200, 200)
line.line.fill.background()

# Itens do timeline
add_timeline_item(slide, 1.5, 2.7, 1, "PILOTO\n4 fornecedores\n701 SKUs", active=True)
add_timeline_item(slide, 4.5, 2.7, 2, "EXPANSÃO\n+8 fornecedores\n~2.000 SKUs", active=False)
add_timeline_item(slide, 7.5, 2.7, 3, "EXPANSÃO\n+11 fornecedores\n~5.000 SKUs", active=False)
add_timeline_item(slide, 10.5, 2.7, 4, "COMPLETO\n23 categorias\nTodos SKUs", active=False)

# Critérios
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12), Inches(1.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Critérios de Expansão"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = AZUL_ESCURO
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = "Resultados positivos na fase anterior  •  Dados históricos disponíveis  •  Capacidade da equipe"
p.font.size = Pt(14)
p.font.color.rgb = CINZA
p.alignment = PP_ALIGN.CENTER

# SLIDE 9 - Cronograma (visual simples)
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "CRONOGRAMA DO PILOTO")

fases = [
    ("Semana 1-2", "SETUP", "Configuração inicial", AZUL_CLARO, 2),
    ("Semana 3-4", "TREINAMENTO", "Capacitação equipe", VERDE, 2),
    ("Semana 5-8", "OPERAÇÃO ASSISTIDA", "Acompanhamento próximo", LARANJA, 4),
    ("Semana 9-12", "OPERAÇÃO AUTÔNOMA", "Avaliação de resultados", RGBColor(128, 0, 128), 4)
]

y = 1.8
for periodo, titulo, desc, cor, semanas in fases:
    # Barra de Gantt simplificada
    largura = semanas * 0.8
    x_start = 4 + (int(periodo.split()[1].split('-')[0]) - 1) * 0.66

    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_start), Inches(y), Inches(largura), Inches(0.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = cor
    bar.line.fill.background()

    # Período
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(y + 0.15), Inches(2), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = periodo
    p.font.size = Pt(14)
    p.font.color.rgb = CINZA

    # Título na barra
    txBox = slide.shapes.add_textbox(Inches(x_start), Inches(y + 0.2), Inches(largura), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER

    y += 1.2

# Marcos
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(12), Inches(0.8))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Marcos: Sem 4 - Primeira previsão  |  Sem 8 - Ciclo completo  |  Sem 12 - Avaliação final"
p.font.size = Pt(14)
p.font.color.rgb = AZUL_ESCURO
p.alignment = PP_ALIGN.CENTER

# SLIDE 10 - Por que agora
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "POR QUE AGORA?")

razoes = [
    ("DADOS", "Histórico carregado\ne validado", AZUL_CLARO),
    ("SISTEMA", "Pronto para\nuso imediato", VERDE),
    ("INVESTIMENTO", "Baixo custo\ninicial", LARANJA),
    ("RETORNO", "Alto potencial\nde ganhos", RGBColor(128, 0, 128)),
    ("RISCO", "Controlado\nvia piloto", VERMELHO)
]

for i, (titulo, desc, cor) in enumerate(razoes):
    x = 0.3 + i * 2.6

    # Círculo
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.65), Inches(2), Inches(1.3), Inches(1.3))
    circle.fill.solid()
    circle.fill.fore_color.rgb = cor
    circle.line.fill.background()

    # Título no círculo
    txBox = slide.shapes.add_textbox(Inches(x + 0.65), Inches(2.4), Inches(1.3), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER

    # Descrição abaixo
    txBox = slide.shapes.add_textbox(Inches(x), Inches(3.5), Inches(2.6), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(14)
    p.font.color.rgb = PRETO
    p.alignment = PP_ALIGN.CENTER

# Call to action
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(5), Inches(7.3), Inches(1.5))
box.fill.solid()
box.fill.fore_color.rgb = AZUL_ESCURO
box.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(3), Inches(5.3), Inches(7.3), Inches(1))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "O MOMENTO É AGORA"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "Vamos começar o piloto?"
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(200, 200, 200)
p.alignment = PP_ALIGN.CENTER

# SLIDE 11 - Final
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Fundo
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
shape.fill.solid()
shape.fill.fore_color.rgb = AZUL_ESCURO
shape.line.fill.background()

# Linha decorativa
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3), Inches(2.3), Inches(0.05))
line.fill.solid()
line.fill.fore_color.rgb = LARANJA
line.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(12.333), Inches(1.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "OBRIGADO"
p.font.size = Pt(56)
p.font.bold = True
p.font.color.rgb = BRANCO
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = "Dúvidas?"
p.font.size = Pt(24)
p.font.color.rgb = RGBColor(200, 200, 200)
p.alignment = PP_ALIGN.CENTER
p.space_before = Pt(30)

# Salvar
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'Apresentacao_Previsao_Demanda_v2.pptx')
prs.save(output_path)
print(f'Apresentação salva em: {output_path}')
