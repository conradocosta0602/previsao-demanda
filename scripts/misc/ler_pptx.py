# -*- coding: utf-8 -*-
"""Script para ler conteudo do PowerPoint"""
from pptx import Presentation
from pptx.util import Inches, Pt

pptx_path = r"C:\temp\apresentacao_fc_v1.pptx"

prs = Presentation(pptx_path)

print(f"Total de slides: {len(prs.slides)}")
print(f"Dimensoes: {prs.slide_width.inches} x {prs.slide_height.inches} inches")
print("\n" + "="*80)

for slide_num, slide in enumerate(prs.slides, 1):
    print(f"\n=== SLIDE {slide_num} ===")
    print(f"Numero de shapes: {len(slide.shapes)}")

    for shape in slide.shapes:
        print(f"\n  Shape: {shape.shape_type}, Pos: ({shape.left.inches:.2f}, {shape.top.inches:.2f}), Size: ({shape.width.inches:.2f}x{shape.height.inches:.2f})")

        if hasattr(shape, "text") and shape.text:
            text_preview = shape.text[:100].replace('\n', ' | ')
            print(f"    Texto: {text_preview}...")

        if shape.has_table:
            table = shape.table
            print(f"    TABELA: {len(table.rows)} linhas x {len(table.columns)} colunas")
            for row_idx, row in enumerate(table.rows):
                row_text = []
                for cell in row.cells:
                    row_text.append(cell.text[:30] if cell.text else "")
                print(f"      Linha {row_idx}: {row_text}")
