# -*- coding: utf-8 -*-
"""
Script para ajustar cronograma do slide 11 de 8 meses para 12 meses
Mantendo as mesmas 4 ações
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import copy

# Abrir apresentação
pptx_input = r"C:\temp\apresentacao_fc_v1.pptx"
pptx_output = r"C:\temp\apresentacao_fc_v2.pptx"

prs = Presentation(pptx_input)

# Pegar o slide 11 (índice 10)
slide = prs.slides[10]

# Identificar shapes para remover (cronograma atual)
shapes_to_remove = []

# Vamos recriar o slide 11 do zero, removendo os elementos do cronograma
# mas mantendo o header

print("Analisando shapes do slide 11...")
for shape in slide.shapes:
    # Identificar elementos do cronograma pelos textos e posições
    if hasattr(shape, "text"):
        text = shape.text
        # Marcadores de mês (M1, M2, etc.)
        if text in ["M1", "M2", "M3", "M4", "M5", "M6", "M7", "M8"]:
            shapes_to_remove.append(shape)
        # Durações
        elif text in ["3 meses", "1 mes", "2 meses"]:
            shapes_to_remove.append(shape)
        # Marcos com "Mes 8"
        elif "Mes 8" in text or "Mes 5" in text or "Mes 3" in text:
            shapes_to_remove.append(shape)

    # Linhas verticais do grid (posição entre 5 e 13 inches)
    if shape.left and shape.width:
        left_inch = shape.left.inches
        width_inch = shape.width.inches
        # Linhas verticais finas
        if width_inch < 0.05 and left_inch > 5:
            shapes_to_remove.append(shape)
        # Barras coloridas do Gantt (entre 5 e 13, e altura ~0.6)
        if left_inch >= 5 and 0.5 < shape.height.inches < 0.7 and shape.width.inches > 0.5:
            shapes_to_remove.append(shape)

print(f"Shapes para remover: {len(shapes_to_remove)}")

# Remover shapes identificados
for shape in shapes_to_remove:
    sp = shape._element
    sp.getparent().remove(sp)

# Agora adicionar o novo cronograma de 12 meses
# Cores
AZUL_ACCENT = RGBColor(0, 120, 200)
VERDE_SUCESSO = RGBColor(40, 167, 69)
LARANJA_ATENCAO = RGBColor(255, 140, 0)
ROXO = RGBColor(128, 90, 180)
AZUL_NAVY = RGBColor(15, 32, 65)
BRANCO = RGBColor(255, 255, 255)
CINZA_CLARO = RGBColor(230, 230, 230)
CINZA_MEDIO = RGBColor(160, 170, 180)

# Posicionamento para 12 meses
# Início X do grid: 4.0 inches
# Largura total disponível: ~9 inches (até 13)
# Cada mês: 9/12 = 0.75 inches

start_x = 4.0
mes_width = 0.75
header_y = 1.7

# Adicionar labels dos meses (M1 a M12)
for i in range(1, 13):
    x = start_x + (i - 1) * mes_width
    txBox = slide.shapes.add_textbox(Inches(x), Inches(header_y), Inches(mes_width), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"M{i}"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = AZUL_NAVY
    p.alignment = PP_ALIGN.CENTER

    # Linha vertical
    if i <= 12:
        vline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + mes_width), Inches(2.0), Inches(0.01), Inches(4.0))
        vline.fill.solid()
        vline.fill.fore_color.rgb = CINZA_CLARO
        vline.line.fill.background()

# Fases do Gantt com novo dimensionamento para 12 meses:
# Setup e Configuração: 4 meses (M1-M4) - era 3
# Treinamento: 1 mes (M5) - mantém
# Operação Assistida: 3 meses (M6-M8) - era 2
# Operação Autônoma: 4 meses (M9-M12) - era 2

fases_gantt = [
    ("Setup e Configuracao", 0, 4, AZUL_ACCENT, "4 meses", 2.2),
    ("Treinamento da Equipe", 4, 1, VERDE_SUCESSO, "1 mes", 3.3),
    ("Operacao Assistida", 5, 3, LARANJA_ATENCAO, "3 meses", 4.4),
    ("Operacao Autonoma", 8, 4, ROXO, "4 meses", 5.5)
]

for fase, start, duration, color, duracao_texto, y_pos in fases_gantt:
    # Barra do Gantt
    bar_x = start_x + start * mes_width
    bar_width = duration * mes_width

    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(bar_x), Inches(y_pos), Inches(bar_width), Inches(0.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    # Texto na barra
    txBox = slide.shapes.add_textbox(Inches(bar_x), Inches(y_pos + 0.15), Inches(bar_width), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = duracao_texto
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER

# Atualizar o box de marcos no rodapé
# Encontrar e atualizar o texto dos marcos
for shape in slide.shapes:
    if hasattr(shape, "text") and "MARCOS:" in shape.text:
        # Atualizar texto
        tf = shape.text_frame
        for paragraph in tf.paragraphs:
            if "MARCOS:" in paragraph.text:
                paragraph.text = "MARCOS:  Mes 4 - Sistema validado  |  Mes 7 - Primeiro ciclo completo  |  Mes 12 - Avaliacao e decisao"

# Salvar
prs.save(pptx_output)
print(f"\nApresentacao salva em: {pptx_output}")
print("Cronograma ajustado de 8 para 12 meses!")
print("\nNovas durações:")
print("  - Setup e Configuração: 4 meses (M1-M4)")
print("  - Treinamento da Equipe: 1 mês (M5)")
print("  - Operação Assistida: 3 meses (M6-M8)")
print("  - Operação Autônoma: 4 meses (M9-M12)")
